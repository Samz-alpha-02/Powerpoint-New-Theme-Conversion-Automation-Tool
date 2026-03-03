"""
template_applier.py
Core logic for applying a new PowerPoint template to existing presentations.

Strategy (v4 – in-place, single-package)
-----------------------------------------
The template may contain decorative shapes (coloured rectangles, logos, lines)
placed directly on the slides – not only in the slide master.

Rather than copying template shapes across OPC packages (which risks dangling
references and repair dialogs), we work directly on the template slides that
are already inside dst_prs:

  1. Load dst_prs from the template (template slides + master/layouts intact).
  2. Ensure dst_prs has exactly N slides (N = source slide count):
       • If template has fewer slides → duplicate last template slide *within*
         dst_prs (same OPC package, reuse existing Part objects, no cross-
         package rId issues).
       • If template has more slides → remove extras.
  3. For each slide pair (dst_slide ↔ src_slide):
       a. Copy source media blobs into dst_prs (blob copy, new Part objects).
       b. Clear placeholder content in the template slide.
       c. Inject source placeholder content matched by idx, then by ph-type.
  4. Apply theme fonts and enable word-wrap / normalAutofit.

All shapes that produce the template's visual design are already embedded in
the template slides – no shape-tree copying between packages is required.
"""

import copy
import logging
from io import BytesIO
from typing import Callable, Optional

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.packuri import PackURI

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Media relationship keywords
# ---------------------------------------------------------------------------
_MEDIA_KWDS = (
    "image", "chart", "diagram", "oleObject", "ole", "audio", "video", "media",
)

# ---------------------------------------------------------------------------
# Theme font helpers
# ---------------------------------------------------------------------------

def _get_theme_fonts(prs: Presentation) -> dict:
    fonts: dict = {"major_font": None, "minor_font": None}
    try:
        sm_part = prs.slide_master.part
        for rel in sm_part.rels.values():
            if "theme" in (rel.reltype or "").lower():
                try:
                    root = etree.fromstring(rel._target.blob)       # type: ignore[attr-defined]
                    fs   = root.find(".//" + qn("a:fontScheme"))
                    if fs is not None:
                        mj = fs.find(qn("a:majorFont") + "/" + qn("a:latin"))
                        mn = fs.find(qn("a:minorFont") + "/" + qn("a:latin"))
                        if mj is not None: fonts["major_font"] = mj.get("typeface")
                        if mn is not None: fonts["minor_font"] = mn.get("typeface")
                except Exception:
                    pass
    except Exception as exc:
        logger.debug("Font extraction error: %s", exc)
    return fonts

# ---------------------------------------------------------------------------
# Layout / template-slide matching
# ---------------------------------------------------------------------------

def _layout_name(obj) -> str:
    try:   return (obj.name or "").strip().lower()
    except Exception: return ""


def _best_layout(template_prs: Presentation, old_name: str):
    layouts = template_prs.slide_master.slide_layouts
    lo = old_name.strip().lower()
    for lay in layouts:
        if _layout_name(lay) == lo: return lay
    for lay in layouts:
        n = _layout_name(lay)
        if lo and (lo in n or n in lo): return lay
    KEYWORDS = {
        "title":   ["title slide", "title only", "title"],
        "content": ["title and content", "two content", "content"],
        "blank":   ["blank"],
        "section": ["section header"],
        "picture": ["picture with caption"],
    }
    for kw, cands in KEYWORDS.items():
        if kw in lo:
            for cand in cands:
                for lay in layouts:
                    if cand in _layout_name(lay): return lay
    return layouts[0]


# ---------------------------------------------------------------------------
# OPC helpers
# ---------------------------------------------------------------------------

def _unique_partname(package, base: str) -> PackURI:
    existing = {p.partname for p in package.iter_parts()}
    if base not in existing:
        return PackURI(base)
    stem, _, ext = base.rpartition(".")
    i = 2
    while True:
        cand = f"{stem}_{i}.{ext}" if ext else f"{stem}_{i}"
        if cand not in existing:
            return PackURI(cand)
        i += 1


def _copy_slide_media(src_part, dst_part) -> dict:
    """
    Copy media / hyperlink rels from src_part → dst_part (possibly different
    packages).  Returns {old_rId: new_rId}.
    """
    from pptx.opc.package import Part

    rId_map: dict = {}
    dst_pkg = dst_part.package

    for old_rId, rel in list(src_part.rels.items()):
        try:
            reltype = rel.reltype or ""
            if rel.is_external:
                if "hyperlink" in reltype.lower():
                    rId_map[old_rId] = dst_part.relate_to(
                        rel.target_ref, reltype, is_external=True
                    )
                continue
            if not any(kw in reltype for kw in _MEDIA_KWDS):
                continue
            tgt          = rel._target                              # type: ignore[attr-defined]
            new_partname = _unique_partname(dst_pkg, tgt.partname)
            new_part     = Part(new_partname, tgt.content_type, dst_pkg, tgt.blob)  # type: ignore[attr-defined]
            rId_map[old_rId] = dst_part.relate_to(new_part, reltype)
        except Exception as exc:
            logger.debug("Skipping rel %s: %s", old_rId, exc)
    return rId_map


def _copy_slide_media_intra(src_part, dst_part) -> dict:
    """
    Copy media / hyperlink rels from src_part → dst_part when **both parts
    live in the same OPC package** (e.g. duplicating a slide within dst_prs).
    We reuse the existing Part objects – just create new relationships.
    Returns {old_rId: new_rId}.
    """
    rId_map: dict = {}
    for old_rId, rel in list(src_part.rels.items()):
        try:
            reltype = rel.reltype or ""
            if rel.is_external:
                if "hyperlink" in reltype.lower():
                    rId_map[old_rId] = dst_part.relate_to(
                        rel.target_ref, reltype, is_external=True
                    )
                continue
            if not any(kw in reltype for kw in _MEDIA_KWDS):
                continue
            # Reuse the actual same Part object (same package) – no blob copy
            rId_map[old_rId] = dst_part.relate_to(rel._target, reltype)  # type: ignore[attr-defined]
        except Exception as exc:
            logger.debug("Intra-rel skip %s: %s", old_rId, exc)
    return rId_map

# ---------------------------------------------------------------------------
# Slide-level operations (remove / duplicate within same prs)
# ---------------------------------------------------------------------------

def _remove_slide(prs: Presentation, index: int) -> None:
    """Remove the slide at *index* from *prs* (0-based)."""
    sld_reltype = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
    )
    slide      = prs.slides[index]
    slide_part = slide.part
    rels_obj   = prs.part._rels
    rid_to_remove = None
    for rId, rel in list(rels_obj._rels.items()):
        if getattr(rel, "reltype", "") == sld_reltype and rel._target is slide_part:   # type: ignore[attr-defined]
            rid_to_remove = rId
            break
    if rid_to_remove:
        try:
            rels_obj.pop(rid_to_remove)
        except Exception:
            rels_obj._rels.pop(rid_to_remove, None)
    sldId_elems = prs.slides._sldIdLst.findall(qn("p:sldId"))
    if index < len(sldId_elems):
        prs.slides._sldIdLst.remove(sldId_elems[index])


def _duplicate_last_slide(prs: Presentation) -> None:
    """
    Append a copy of the last slide in *prs* using only parts already inside
    the same OPC package (no cross-package blob copies).
    """
    src_slide  = prs.slides[-1]
    src_layout = src_slide.slide_layout
    new_slide  = prs.slides.add_slide(src_layout)

    src_sp_tree = src_slide.shapes._spTree
    dst_sp_tree = new_slide.shapes._spTree
    for child in list(dst_sp_tree)[2:]:
        dst_sp_tree.remove(child)

    rId_map = _copy_slide_media_intra(src_slide.part, new_slide.part)

    for child in list(src_sp_tree)[2:]:
        cloned = copy.deepcopy(child)
        _remap_rids(cloned, rId_map)
        dst_sp_tree.append(cloned)

    try:
        src_bg = src_slide.element.find(qn("p:bg"))
        if src_bg is not None:
            dst_cSld = new_slide.element.find(qn("p:cSld"))
            existing_bg = dst_cSld.find(qn("p:bg"))
            if existing_bg is not None:
                dst_cSld.remove(existing_bg)
            new_bg = copy.deepcopy(src_bg)
            _remap_rids(new_bg, rId_map)
            dst_cSld.insert(0, new_bg)
    except Exception:
        pass


def _remap_rids(elem: etree._Element, rId_map: dict) -> None:
    if not rId_map: return
    for el in elem.iter():
        for attr in list(el.attrib):
            v = el.get(attr)
            if v in rId_map: el.set(attr, rId_map[v])

# ---------------------------------------------------------------------------
# Placeholder helpers
# ---------------------------------------------------------------------------

def _ph_idx(shape_elem) -> Optional[int]:
    """Return the placeholder idx attribute of a shape element, or None."""
    try:
        ph = shape_elem.find(".//" + qn("p:ph"))
        if ph is not None:
            idx = ph.get("idx", "0")
            return int(idx)
    except Exception:
        pass
    return None


def _ph_type(shape_elem) -> str:
    """Return the placeholder type string (e.g. 'title', 'body', 'pic')."""
    try:
        ph = shape_elem.find(".//" + qn("p:ph"))
        if ph is not None:
            return ph.get("type", "body")
    except Exception:
        pass
    return ""


def _is_placeholder(shape_elem) -> bool:
    return shape_elem.find(".//" + qn("p:ph")) is not None


def _clear_placeholder_content(sp_elem) -> None:
    """Wipe all text and media from a placeholder, leaving one empty paragraph."""
    txBody = sp_elem.find(qn("p:txBody"))
    if txBody is not None:
        for para in list(txBody.findall(qn("a:p"))):
            txBody.remove(para)
        etree.SubElement(txBody, qn("a:p"))
    # For picture placeholders, remove blipFill so old image doesn't show
    for tag in (qn("p:blipFill"), qn("a:blipFill")):
        bf = sp_elem.find(tag)
        if bf is not None:
            sp_elem.remove(bf)


def _copy_text_content(src_sp, dst_sp) -> None:
    """Copy the txBody paragraphs from src_sp into dst_sp, preserving dst formatting skeleton."""
    src_txBody = src_sp.find(qn("p:txBody"))
    dst_txBody = dst_sp.find(qn("p:txBody"))
    if src_txBody is None or dst_txBody is None:
        return
    # Remove existing paragraphs from dst
    for p in list(dst_txBody.findall(qn("a:p"))):
        dst_txBody.remove(p)
    # Copy paragraphs from src
    for p in src_txBody.findall(qn("a:p")):
        dst_txBody.append(copy.deepcopy(p))


def _copy_pic_content(src_sp, dst_sp, rId_map: dict) -> None:
    """Copy picture blipFill rId from src into dst and remap."""
    src_blip = src_sp.find(".//" + qn("a:blip"))
    dst_blip = dst_sp.find(".//" + qn("a:blip"))
    if src_blip is None:
        return
    if dst_blip is None:
        # Try to copy the whole blipFill
        src_bf = src_sp.find(".//" + qn("p:blipFill"))
        if src_bf is None:
            src_bf = src_sp.find(".//" + qn("a:blipFill"))
        if src_bf is not None:
            dst_sp.append(copy.deepcopy(src_bf))
            _remap_rids(dst_sp, rId_map)
        return
    embed = src_blip.get(qn("r:embed"))
    if embed and embed in rId_map:
        dst_blip.set(qn("r:embed"), rId_map[embed])

# ---------------------------------------------------------------------------
# Strip run-level text colour overrides (so theme colours apply)
# ---------------------------------------------------------------------------

_FILL_TAGS = frozenset({
    qn("a:solidFill"), qn("a:gradFill"), qn("a:blipFill"),
    qn("a:pattFill"),  qn("a:grpFill"),
})

def _strip_text_colours(sp_tree: etree._Element) -> None:
    """Remove explicit text run colour fills so theme colours show through."""
    for rpr in sp_tree.iter(qn("a:rPr"), qn("a:endParaRPr"), qn("a:defRPr")):
        for child in list(rpr):
            if child.tag in _FILL_TAGS:
                rpr.remove(child)
    for ppr in sp_tree.iter(qn("a:pPr")):
        for child in list(ppr):
            if child.tag == qn("a:buClr"):
                ppr.remove(child)

# ---------------------------------------------------------------------------
# Ensure layout-inherited placeholders are materialised in the slide XML
# ---------------------------------------------------------------------------

def _ensure_slide_placeholders(slide) -> None:
    """
    When a slide inherits a placeholder from its layout without overriding it,
    no <p:sp> element exists in the slide's own spTree.  We can't inject
    content into a non-existent element, so we clone a minimal blank override
    from the layout for every ph-idx that is missing at the slide level.
    """
    sp_tree = slide.shapes._spTree
    existing_idx: set = set()
    for sp in list(sp_tree):
        ph = sp.find(".//" + qn("p:ph"))
        if ph is not None:
            existing_idx.add(int(ph.get("idx", "0")))

    # ph types that carry meta-info (date/footer/slide-number) – not content
    _META_TYPES = {"dt", "ftr", "sldNum"}

    layout = slide.slide_layout
    for layout_sp in list(layout.shapes._spTree):
        if layout_sp.tag != qn("p:sp"):
            continue
        ph = layout_sp.find(".//" + qn("p:ph"))
        if ph is None:
            continue
        # Skip date / footer / slide-number – they are not content slots
        if ph.get("type", "") in _META_TYPES:
            continue
        ph_idx = int(ph.get("idx", "0"))
        if ph_idx in existing_idx:
            continue
        # Clone the layout shape; strip its hint text so nothing leaks through
        new_sp = copy.deepcopy(layout_sp)
        txBody = new_sp.find(qn("p:txBody"))
        if txBody is not None:
            for p in list(txBody.findall(qn("a:p"))):
                txBody.remove(p)
            etree.SubElement(txBody, qn("a:p"))
        sp_tree.append(new_sp)
        existing_idx.add(ph_idx)


# ---------------------------------------------------------------------------
# Textbox injection fallback (for templates with no standard placeholders)
# ---------------------------------------------------------------------------

def _sh_top(shape_elem) -> int:
    """Return the top (y) offset of a shape element in EMUs, for sorting."""
    try:
        spPr = shape_elem.find(qn("p:spPr"))
        if spPr is None:
            return 999999999
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is None:
            return 999999999
        off = xfrm.find(qn("a:off"))
        return int(off.get("y", "999999999")) if off is not None else 999999999
    except Exception:
        return 999999999


def _inject_into_textboxes(src_slide, dst_slide, src_rId_map: dict) -> None:
    """
    When the template slide uses plain TextBoxes (no standard placeholder
    shapes), match source placeholder content to the template TextBoxes by
    vertical position:  topmost TextBox ← source title;
                        next TextBox(es) ← source body/content.

    Only injects if at least one TextBox already holds some text (the template
    sample text), meaning the author intended it as a content slot.
    Only overwrites TextBoxes that currently contain text – decorative boxes
    (filled rectangles, etc.) typically have no visible text content.
    """
    dst_sp_tree = dst_slide.shapes._spTree

    # Collect template TextBoxes (non-placeholder shapes with a text body)
    tpl_tbs = []
    for sp in list(dst_sp_tree)[2:]:
        if _is_placeholder(sp):
            continue
        txBody = sp.find(qn("p:txBody"))
        if txBody is None:
            continue
        # Only include boxes that already have some text (template sample text)
        text = "".join(t.text or "" for t in txBody.iter(qn("a:t")))
        if text.strip():
            tpl_tbs.append(sp)

    if not tpl_tbs:
        return  # nothing to inject into

    # Sort by y-position so title (top) comes first
    tpl_tbs.sort(key=_sh_top)

    # Collect source placeholder content in priority order:
    # title / ctrTitle first, then body / object / other
    src_sp_tree = src_slide.shapes._spTree
    src_title_elems:   list = []
    src_content_elems: list = []

    for sp in list(src_sp_tree)[2:]:
        if not _is_placeholder(sp):
            continue
        ph_type = _ph_type(sp)
        txBody = sp.find(qn("p:txBody"))
        has_text = txBody is not None and any(
            t.text and t.text.strip() for t in txBody.iter(qn("a:t"))
        )
        if not has_text:
            continue
        if ph_type in ("title", "ctrTitle"):
            src_title_elems.append(sp)
        else:
            src_content_elems.append(sp)

    src_ordered = src_title_elems + src_content_elems

    # Inject: template TextBox[0] ← src[0] (title), TextBox[1] ← src[1] …
    # First clear ALL template TextBoxes so sample text never leaks through
    for tpl_tb in tpl_tbs:
        txBody = tpl_tb.find(qn("p:txBody"))
        if txBody is not None:
            for p in list(txBody.findall(qn("a:p"))):
                txBody.remove(p)
            etree.SubElement(txBody, qn("a:p"))

    for tpl_tb, src_sp in zip(tpl_tbs, src_ordered):
        _copy_text_content(src_sp, tpl_tb)
        _remap_rids(tpl_tb, src_rId_map)

    _strip_text_colours(dst_sp_tree)


# ---------------------------------------------------------------------------
# Core: inject source content into an existing template slide
# ---------------------------------------------------------------------------

def _inject_content(src_slide, dst_slide) -> None:
    """
    Inject source slide placeholder content into *dst_slide* (a template slide
    already present in dst_prs).

    Steps
    -----
    1. Copy source media into dst_slide's part.
    2. Clear all placeholder content in dst_slide.
    3. Build idx → element and type → element lookups for dst placeholders.
    4. For each source placeholder, find matching dst placeholder and copy
       text / image content.
    5. Strip inline colour overrides so the new theme colours show through.

    Source non-placeholder shapes are intentionally NOT appended – they are
    decoration from the OLD template and would appear at wrong positions on
    the new canvas.
    """
    # ── 1. Copy source media ───────────────────────────────────────────────
    src_rId_map = _copy_slide_media(src_slide.part, dst_slide.part)

    # Materialise any layout-inherited placeholders so we can inject into them
    _ensure_slide_placeholders(dst_slide)

    dst_sp_tree = dst_slide.shapes._spTree
    src_sp_tree = src_slide.shapes._spTree

    # ── 2. Clear template placeholder content ─────────────────────────────
    for shape_elem in list(dst_sp_tree)[2:]:
        if _is_placeholder(shape_elem):
            _clear_placeholder_content(shape_elem)

    # ── 3. Build dst lookups (skip meta placeholders: date/footer/sldNum) ──
    _META_PH_TYPES = {"dt", "ftr", "sldNum"}
    dst_ph_by_idx:  dict = {}
    dst_ph_by_type: dict = {}
    for shape_elem in list(dst_sp_tree)[2:]:
        if _is_placeholder(shape_elem):
            idx = _ph_idx(shape_elem)
            typ = _ph_type(shape_elem)
            if typ in _META_PH_TYPES:
                continue   # never use date/footer/slide-number as content slots
            if idx is not None:
                dst_ph_by_idx[idx] = shape_elem
            if typ:
                dst_ph_by_type.setdefault(typ, shape_elem)

    if not dst_ph_by_idx and not dst_ph_by_type:
        # No real content placeholders found – skip straight to textbox injection
        _strip_text_colours(dst_sp_tree)
        _inject_into_textboxes(src_slide, dst_slide, src_rId_map)
        return

    # ── 4. Inject source placeholder content ──────────────────────────────
    for src_shape in list(src_sp_tree)[2:]:
        if not _is_placeholder(src_shape):
            continue  # skip old-template decorations – they'd land off-canvas

        src_idx  = _ph_idx(src_shape)
        src_type = _ph_type(src_shape)

        # Priority: match by index, then by type, then by position
        dst_shape = dst_ph_by_idx.get(src_idx) if src_idx is not None else None
        if dst_shape is None and src_type:
            dst_shape = dst_ph_by_type.get(src_type)
        if dst_shape is None:
            # Last resort: first placeholder for title, second for body
            all_dst = list(dst_ph_by_idx.values())
            if src_type in ("title", "ctrTitle") and all_dst:
                dst_shape = all_dst[0]
            elif all_dst and len(all_dst) > 1:
                dst_shape = all_dst[1]
            elif all_dst:
                dst_shape = all_dst[0]
        if dst_shape is None:
            continue

        # Copy text content
        _copy_text_content(src_shape, dst_shape)

        # Remap hyperlink / media rIds inside copied text
        _remap_rids(dst_shape, src_rId_map)

        # Handle picture placeholders
        is_pic_ph = (
            src_type == "pic"
            or src_shape.find(".//" + qn("p:blipFill")) is not None
            or src_shape.find(".//" + qn("a:blipFill")) is not None
        )
        if is_pic_ph:
            _copy_pic_content(src_shape, dst_shape, src_rId_map)

    # ── 5. Strip text colour overrides ────────────────────────────────────
    _strip_text_colours(dst_sp_tree)

    # ── 6. Fallback: if NO placeholders were injected, try textbox matching ─
    # (happens when the template uses plain TextBoxes instead of Placeholders)
    _inject_into_textboxes(src_slide, dst_slide, src_rId_map)

    # ── 7. Copy slide notes ────────────────────────────────────────────────
    try:
        src_tf = src_slide.notes_slide.notes_placeholder.text_frame
        dst_tf = dst_slide.notes_slide.notes_placeholder.text_frame
        dst_txBody = dst_tf._txBody
        for p in list(dst_txBody):
            dst_txBody.remove(p)
        for p in src_tf._txBody.findall(qn("a:p")):
            dst_txBody.append(copy.deepcopy(p))
    except Exception:
        pass

# ---------------------------------------------------------------------------
# Font & overflow helpers
# ---------------------------------------------------------------------------

def _apply_fonts(slide, theme_fonts: dict) -> None:
    major = theme_fonts.get("major_font")
    minor = theme_fonts.get("minor_font")
    if not major and not minor: return
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        is_title = False
        try:
            if shape.is_placeholder:
                ph_str = str(shape.placeholder_format.type)
                is_title = any(t in ph_str for t in ("TITLE", "CENTER_TITLE"))
        except Exception: pass
        font_name = major if is_title else minor
        if not font_name: continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try: run.font.name = font_name
                except Exception: pass


def _fix_overflow(slide) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        try:
            shape.text_frame.word_wrap = True
            bp = shape.text_frame._txBody.find(qn("a:bodyPr"))
            if bp is not None:
                for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
                    for old in bp.findall(qn(tag)): bp.remove(old)
                etree.SubElement(bp, qn("a:normAutofit"))
        except Exception: pass

# ---------------------------------------------------------------------------
# Public API – single file
# ---------------------------------------------------------------------------

def apply_template(
    old_pptx_bytes: bytes,
    template_pptx_bytes: bytes,
    progress_callback: Optional[Callable[[str], None]] = None,
) -> bytes:
    """Apply the new company template to a single presentation."""

    def _p(msg: str):
        if progress_callback:
            progress_callback(msg)
        logger.info(msg)

    _p("Loading presentations …")
    src_prs = Presentation(BytesIO(old_pptx_bytes))

    # dst_prs is loaded from the template.  Its slides already contain all
    # decorative shapes (rectangles, ovals, logos).  We only manipulate the
    # placeholder *content* – this avoids any cross-package OPC operations
    # for the template shapes themselves.
    dst_prs = Presentation(BytesIO(template_pptx_bytes))

    _p("Extracting theme fonts …")
    theme_fonts = _get_theme_fonts(dst_prs)
    _p(
        f"  Heading: {theme_fonts.get('major_font') or '(from theme)'}  |  "
        f"Body: {theme_fonts.get('minor_font') or '(from theme)'}"
    )

    n_src = len(src_prs.slides)
    n_tpl = len(dst_prs.slides)
    _p(f"Source slides: {n_src}  |  Template slides: {n_tpl}")

    # ── Adjust dst_prs slide count to match source ───────────────────────
    while len(dst_prs.slides) < n_src:
        _p(f"  Duplicating template slide to reach {n_src} slides …")
        _duplicate_last_slide(dst_prs)

    while len(dst_prs.slides) > n_src:
        _p(f"  Removing excess template slide (now {len(dst_prs.slides)} > {n_src}) …")
        _remove_slide(dst_prs, len(dst_prs.slides) - 1)

    _p(f"Processing {n_src} slide(s) …")
    for i, src_slide in enumerate(src_prs.slides):
        _p(f"  Slide {i + 1}/{n_src} …")
        try:
            dst_slide = dst_prs.slides[i]
            _inject_content(src_slide, dst_slide)
        except Exception as exc:
            _p(f"  ⚠ Slide {i + 1} error: {exc}")
            logger.exception("Slide %d injection failed", i + 1)

    _p("Applying theme fonts …")
    for slide in dst_prs.slides:
        _apply_fonts(slide, theme_fonts)

    _p("Fixing text overflow …")
    for slide in dst_prs.slides:
        _fix_overflow(slide)

    _p("Saving …")
    out = BytesIO()
    dst_prs.save(out)
    return out.getvalue()

# ---------------------------------------------------------------------------
# Public API – batch ZIP
# ---------------------------------------------------------------------------

def process_zip(
    zip_bytes: bytes,
    template_bytes: bytes,
    progress_callback: Optional[Callable[[str], None]] = None,
) -> bytes:
    import zipfile

    def _p(msg: str):
        if progress_callback: progress_callback(msg)
        logger.info(msg)

    input_zip  = zipfile.ZipFile(BytesIO(zip_bytes))
    out_buf    = BytesIO()
    output_zip = zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED)

    pptx_names = [
        n for n in input_zip.namelist()
        if n.lower().endswith(".pptx") and not n.startswith("__MACOSX")
    ]
    if not pptx_names:
        raise ValueError(
            "No .pptx files found in the uploaded ZIP. "
            "Please ensure the archive contains at least one .pptx file."
        )

    _p(f"Found {len(pptx_names)} presentation(s) in the ZIP.")

    for idx, name in enumerate(pptx_names, 1):
        _p(f"\n[{idx}/{len(pptx_names)}]  {name}")
        try:
            old_bytes = input_zip.read(name)
            new_bytes = apply_template(
                old_bytes,
                template_bytes,
                progress_callback=lambda m: _p(f"  {m}"),
            )
            output_zip.writestr(name, new_bytes)
            _p("  ✓ Done")
        except Exception as exc:
            _p(f"  ✗ FAILED: {exc}")
            logger.exception("Failed to process %s", name)
            try:   output_zip.writestr(name, input_zip.read(name))
            except Exception: pass

    output_zip.close()
    _p("\nAll presentations processed.")
    return out_buf.getvalue()
