# =============================================================================
# Document Logo Replacer — app.py
# =============================================================================
# A Streamlit web application that replaces logos and text inside PPTX, DOCX,
# PDF, and ZIP (batch) files.  The core strategy is always to operate at the
# raw ZIP / byte level so that none of the document's other content — shapes,
# animations, XML relationships, z-order — is ever touched or re-serialised.
# =============================================================================

import os
import re
import shutil
import tempfile
import zipfile
from io import BytesIO

import imagehash          # perceptual image hashing (pHash) for visual similarity
import streamlit as st
from PIL import Image, ImageDraw   # image manipulation and annotation

# Case-insensitive pattern that matches every capitalisation variant of the
# brand name the user wants to rename (e.g. LTIMindtree / LTIMINDTREE /
# ltimindtree / Ltimindtree / LTIMINDTree …).
_LTIMINDTREE_RE = re.compile(r"LTIMindtree", re.IGNORECASE)


# ============================================================
# BASIC IMAGE HELPERS
# ============================================================

def get_image_hash(img_bytes: bytes):
    """
    Compute a perceptual hash (pHash) of an image.

    Perceptual hashing produces a 64-bit fingerprint that stays similar for
    images that are visually alike (e.g. same logo at slightly different
    resolution or compression).  Two hashes can be compared with `-` to get
    a Hamming distance: 0 = identical, ~10 = very similar, >20 = different.
    """
    return imagehash.phash(Image.open(BytesIO(img_bytes)))


def to_png(img_bytes: bytes) -> bytes:
    """
    Convert any supported image format to a 32-bit RGBA PNG in memory.

    Used when embedding replacement images into PPTX / DOCX ZIP archives,
    ensuring a consistent, lossless format regardless of what the user uploads.
    """
    buf = BytesIO()
    Image.open(BytesIO(img_bytes)).convert("RGBA").save(buf, format="PNG")
    return buf.getvalue()


def shrink_logo(img_bytes: bytes, ratio: float = 0.54) -> bytes:
    """
    Resize a logo to `ratio` of its original dimensions (default 54% = 46% reduction)
    and return the result as an RGBA PNG.

    The resized logo is centred on a transparent canvas the same size as the
    original so that the bounding-box slot in the document is filled correctly
    and the surrounding area remains transparent (for PPTX/DOCX) or white (for PDF).
    """
    img = Image.open(BytesIO(img_bytes)).convert("RGBA")
    orig_w, orig_h = img.size
    new_w = max(1, int(orig_w * ratio))
    new_h = max(1, int(orig_h * ratio))
    resized = img.resize((new_w, new_h), Image.LANCZOS)
    # Place the smaller logo centred on a transparent canvas of the original size
    canvas = Image.new("RGBA", (orig_w, orig_h), (0, 0, 0, 0))
    offset_x = (orig_w - new_w) // 2
    offset_y = (orig_h - new_h) // 2
    canvas.paste(resized, (offset_x, offset_y), resized)
    buf = BytesIO()
    canvas.save(buf, format="PNG")
    return buf.getvalue()


def thumbnail_bytes(img_bytes: bytes, size: tuple = (220, 220)) -> bytes:
    """
    Create a small RGBA thumbnail for gallery display in the Streamlit UI.

    Uses `Image.thumbnail` which preserves aspect ratio and never upscales.
    LANCZOS resampling gives the best quality when downscaling.
    """
    img = Image.open(BytesIO(img_bytes)).convert("RGBA")
    img.thumbnail(size, Image.LANCZOS)
    buf = BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def annotate_bbox(
    img_bytes: bytes,
    x1: int, y1: int, x2: int, y2: int,
    color: str = "#ff3300",
    line_width: int = 4,
) -> bytes:
    """
    Draw a visible bounding-box rectangle onto an image for preview purposes.

    The rectangle is drawn `line_width` times, each iteration offset by 1 px
    outward, to produce a thick visible border without using PIL's `width`
    parameter (which is unavailable on older Pillow builds).

    Returns a PNG byte-string suitable for `st.image()`.
    """
    img = Image.open(BytesIO(img_bytes)).convert("RGBA")
    draw = ImageDraw.Draw(img)
    # Draw concentric rectangles expanding outward to simulate line thickness
    for i in range(line_width):
        draw.rectangle([x1 - i, y1 - i, x2 + i, y2 + i], outline=color)
    buf = BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ============================================================
# OCR  (EasyOCR, cached so model loads once)
# ============================================================

@st.cache_resource(show_spinner="Loading OCR engine…")
def _ocr_reader():
    """
    Initialise and cache the EasyOCR English reader.

    `@st.cache_resource` ensures the ~300 MB language model is downloaded and
    loaded into memory only once per Streamlit session, no matter how many
    times the user triggers detection.  `gpu=False` keeps it CPU-only for
    broad compatibility.
    """
    import easyocr  # lazy import — only needed when OCR mode is used
    return easyocr.Reader(["en"], gpu=False)


def detect_by_ocr(
    img_bytes: bytes,
    search_text: str = "LTIMindtree",
    expand_left_ratio: float = 0.9,
    v_padding_ratio: float = 0.15,
) -> dict | None:
    """
    Locate a text string inside an image using EasyOCR and return its bounding box.

    Because many logos place a graphical icon *immediately to the left* of the
    text (e.g. the LTIMindtree circular symbol), the detected text bbox is
    expanded leftward by `expand_left_ratio × text_width` so the full logo
    (icon + text) is captured as one region.

    Parameters
    ----------
    img_bytes         : Raw image bytes (any PIL-readable format).
    search_text       : The string to look for (case-insensitive substring match).
    expand_left_ratio : How far left to extend the box relative to text width.
                        0.9 = expand 90 % of the text width to the left.
    v_padding_ratio   : Vertical padding added top & bottom as a fraction of
                        text height, to avoid clipping descenders.

    Returns
    -------
    dict with keys x1, y1, x2, y2, text, confidence  — or None if not found.
    """
    import numpy as np

    reader  = _ocr_reader()
    img     = Image.open(BytesIO(img_bytes)).convert("RGB")
    # readtext returns list of (bbox_polygon, text_string, confidence_float)
    results = reader.readtext(np.array(img))

    best = None
    for bbox, text, conf in results:
        if search_text.lower() in text.lower():
            # bbox is a list of 4 (x,y) corner points; compute axis-aligned rect
            xs = [p[0] for p in bbox]
            ys = [p[1] for p in bbox]
            rx1, ry1 = int(min(xs)), int(min(ys))
            rx2, ry2 = int(max(xs)), int(max(ys))
            # Keep only the highest-confidence match when multiple runs match
            if best is None or conf > best["confidence"]:
                best = {"x1": rx1, "y1": ry1, "x2": rx2, "y2": ry2,
                        "text": text, "confidence": conf}

    if best is None:
        return None

    # Expand bounding box to capture the icon to the left of the text
    w = best["x2"] - best["x1"]
    h = best["y2"] - best["y1"]
    pad_v    = int(h * v_padding_ratio)
    pad_left = int(w * expand_left_ratio)

    img_w, img_h = img.size
    # Clamp so the box never exceeds image boundaries
    best["x1"] = max(0, best["x1"] - pad_left)
    best["y1"] = max(0, best["y1"] - pad_v)
    best["y2"] = min(img_h, best["y2"] + pad_v)
    return best


# ============================================================
# SCREENSHOT PREPROCESSING
# ============================================================

def autocrop_screenshot(img_bytes: bytes, tolerance: int = 30) -> bytes:
    """
    Strip uniform solid-colour borders from a screenshot before template matching.

    When a user captures a screenshot, it often includes desktop wallpaper,
    window chrome, or application padding around the actual logo.  Those extra
    pixels confuse template matching because they don't appear in the embedded
    document image.

    Algorithm:
      1. Sample the top-left corner pixel as the assumed background colour.
      2. Walk inward from each edge, removing rows/columns whose pixels are all
         within `tolerance` (Euclidean per-channel difference) of that colour.
      3. Return the cropped image, or the original if nothing was trimmed.

    `tolerance=30` handles slight JPEG compression artefacts at the border.
    """
    import numpy as np

    img  = Image.open(BytesIO(img_bytes)).convert("RGB")
    arr  = np.array(img)
    h, w = arr.shape[:2]

    # Use top-left corner pixel as background colour reference
    bg = arr[0, 0].astype(int)

    # Helper lambdas: True when every pixel in a row/col is within tolerance of bg
    def _is_border_row(row):  return np.all(np.abs(arr[row].astype(int) - bg) <= tolerance)
    def _is_border_col(col):  return np.all(np.abs(arr[:, col].astype(int) - bg) <= tolerance)

    # Advance each edge inward until a non-background row/column is found
    top    = 0
    while top < h and _is_border_row(top):          top    += 1
    bottom = h - 1
    while bottom > top and _is_border_row(bottom):  bottom -= 1
    left   = 0
    while left < w and _is_border_col(left):        left   += 1
    right  = w - 1
    while right > left and _is_border_col(right):   right  -= 1

    if top >= bottom or left >= right:   # entire image was background — return as-is
        return img_bytes

    cropped = img.crop((left, top, right + 1, bottom + 1))
    buf = BytesIO()
    cropped.save(buf, format="PNG")
    return buf.getvalue()


def _normalise_for_matching(gray_img):
    """
    Apply CLAHE (Contrast Limited Adaptive Histogram Equalisation) to a
    grayscale image before template matching.

    CLAHE redistributes pixel intensities locally so that a screenshot taken
    on a bright monitor matches an image embedded in a document at a different
    brightness level.  `clipLimit=2.0` prevents over-amplification of noise,
    and `tileGridSize=(8,8)` divides the image into 64 tiles for local
    normalisation.
    """
    import cv2
    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
    return clahe.apply(gray_img)


# ============================================================
# TEMPLATE MATCHING  (OpenCV multi-scale)
# ============================================================

def detect_by_template(
    img_bytes: bytes,
    template_bytes: bytes,
    min_score: float = 0.45,
) -> dict | None:
    """
    Locate a logo inside an image using multi-scale OpenCV template matching.

    Why multi-scale?
    The logo in a document may be rendered at a different pixel size than the
    screenshot the user uploads as a reference.  We try scales from 40 % to
    160 % of the template dimensions (in 5 % steps) and keep the best match.

    Pipeline:
      1. Autocrop the user's screenshot to remove desktop padding.
      2. Convert both images to grayscale and apply CLAHE normalisation so
         brightness differences don't hurt matching quality.
      3. At each scale, resize the template and run `cv2.matchTemplate` with
         the TM_CCOEFF_NORMED metric (score 0–1, 1 = perfect match).
      4. Return the best bounding box if its score exceeds `min_score`.

    Parameters
    ----------
    img_bytes      : The full image to search within (e.g. a slide background).
    template_bytes : The reference screenshot/crop of the logo.
    min_score      : Minimum normalised correlation score to count as a match.
                     0.45 works well for high-quality logos; lower values allow
                     more tolerance for heavily compressed images.

    Returns
    -------
    dict with x1, y1, x2, y2, score  — or None if no match above min_score.
    """
    import cv2
    import numpy as np

    # Remove padding/background from user screenshot before matching
    template_bytes = autocrop_screenshot(template_bytes)

    # Convert both to normalised grayscale for invariance to colour/brightness
    img  = _normalise_for_matching(
        cv2.cvtColor(np.array(Image.open(BytesIO(img_bytes)).convert("RGB")),      cv2.COLOR_RGB2GRAY)
    )
    tmpl = _normalise_for_matching(
        cv2.cvtColor(np.array(Image.open(BytesIO(template_bytes)).convert("RGB")), cv2.COLOR_RGB2GRAY)
    )
    ih, iw = img.shape
    th, tw = tmpl.shape

    best_score, best_loc, best_scale = -1.0, None, 1.0

    # Iterate over scales 40 %–160 % in 5 % increments
    for scale in [s / 100 for s in range(40, 165, 5)]:
        new_w = max(1, int(tw * scale))
        new_h = max(1, int(th * scale))
        if new_w > iw or new_h > ih:
            continue   # scaled template larger than image — skip
        resized = cv2.resize(tmpl, (new_w, new_h), interpolation=cv2.INTER_AREA)
        res = cv2.matchTemplate(img, resized, cv2.TM_CCOEFF_NORMED)
        _, score, _, loc = cv2.minMaxLoc(res)
        if score > best_score:
            best_score, best_loc, best_scale = score, loc, scale

    if best_score < min_score or best_loc is None:
        return None   # no acceptable match found

    # Convert best match back to pixel coordinates in the original image
    bw = int(tw * best_scale)
    bh = int(th * best_scale)
    x1, y1 = best_loc
    return {"x1": x1, "y1": y1, "x2": x1 + bw, "y2": y1 + bh, "score": float(best_score)}


# ============================================================
# REGION REPLACE  (paste new logo into detected bbox)
# ============================================================

def replace_region_in_image(
    img_bytes: bytes,
    x1: int, y1: int, x2: int, y2: int,
    new_logo_bytes: bytes,
) -> bytes:
    """
    Paint the new logo over a specific rectangular region of an existing image.

    Used when a logo is *baked into* a background image (not a standalone shape)
    and cannot be replaced by swapping media files.  The detected bounding box
    (from OCR or template matching) defines where to paint.

    Steps:
      1. Open the host image as RGBA.
      2. Calculate the exact pixel dimensions of the bounding box.
      3. Resize the new logo to exactly fit the bounding box (preserving aspect
         ratio is intentionally NOT done here — the new logo fills the space
         that the old logo occupied).
      4. Fill the bounding box with opaque white first to erase any remnants of
         the old logo before compositing the new one.
      5. Alpha-composite the new logo on top.

    Returns the modified image as a PNG byte-string.
    """
    img      = Image.open(BytesIO(img_bytes)).convert("RGBA")
    new_w    = max(1, x2 - x1)
    new_h    = max(1, y2 - y1)
    new_logo = Image.open(BytesIO(new_logo_bytes)).convert("RGBA")
    new_logo = new_logo.resize((new_w, new_h), Image.LANCZOS)

    # Erase the old logo area with solid white before pasting the new one
    clear = Image.new("RGBA", (new_w, new_h), (255, 255, 255, 255))
    img.paste(clear, (x1, y1))
    # Composite new logo using its own alpha channel as the mask
    img.paste(new_logo, (x1, y1), new_logo)

    buf = BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ============================================================
# DOCUMENT EXTRACTION  (returns zip_path / xref for re-embedding)
# ============================================================

def _collect_pics(shapes):
    """
    Recursively collect all PICTURE shapes from a python-pptx shape tree.

    Shapes can be nested inside GROUP shapes, so this recurses into any group.
    Returns a flat list of all picture shape objects found at any depth.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    pics = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Recursively descend into group shapes
            pics.extend(_collect_pics(shape.shapes))
    return pics


def extract_images_pptx(pptx_bytes: bytes) -> list[dict]:
    """
    Extract all unique media images from a PPTX file with their locations.

    Two-phase approach:
      Phase 1 — ZIP scan:  Read every file under `ppt/media/` directly from
        the ZIP to get reliable zip_path values needed for later byte-swapping.
        Images are de-duplicated by perceptual hash so the user sees each
        unique logo only once even if it appears on multiple slides.

      Phase 2 — python-pptx annotation:  Walk the slide master, layouts, and
        individual slides to map each image hash back to human-readable
        location names (e.g. "Slide Master", "Layout — Title Slide",
        "Slide 3").

    The result list is displayed in the Streamlit gallery for the user to
    select which images are logos.

    Returns a list of dicts with keys:
      bytes     — raw image bytes
      zip_path  — internal ZIP path used for replacement (e.g. ppt/media/image1.png)
      locations — list of location strings
      hash_str  — perceptual hash string (used as de-dup key)
    """
    from pptx import Presentation

    # ── Phase 1: collect all media files with their zip paths ──────────────
    media: dict[str, dict] = {}   # hash_str → entry
    with zipfile.ZipFile(BytesIO(pptx_bytes)) as z:
        for name in z.namelist():
            if name.lower().startswith("ppt/media/"):
                try:
                    data = z.read(name)
                    key  = str(get_image_hash(data))
                    if key not in media:   # first occurrence wins for zip_path
                        media[key] = {
                            "bytes":    data,
                            "zip_path": name,
                            "locations": [],
                            "hash_str": key,
                        }
                except Exception:
                    pass  # skip corrupted or unreadable media entries

    # ── Phase 2: annotate with human-readable location names ───────────────
    prs = Presentation(BytesIO(pptx_bytes))

    def annotate(shapes, location: str):
        """Add `location` to the media entry for every picture shape found."""
        for pic in _collect_pics(shapes):
            try:
                key = str(get_image_hash(pic.image.blob))
                if key in media and location not in media[key]["locations"]:
                    media[key]["locations"].append(location)
            except Exception:
                pass

    annotate(prs.slide_master.shapes, "Slide Master")
    for layout in prs.slide_master.slide_layouts:
        annotate(layout.shapes, f"Layout — {layout.name}")
    for i, slide in enumerate(prs.slides):
        annotate(slide.shapes, f"Slide {i + 1}")

    # Any media file not reachable through a shape gets a generic label
    for entry in media.values():
        if not entry["locations"]:
            entry["locations"].append("Media file (not linked to a shape)")

    return list(media.values())


def extract_images_docx(docx_bytes: bytes) -> list[dict]:
    """
    Extract all unique embedded images from a DOCX file.

    DOCX is also a ZIP archive; images live under `word/media/`.
    De-duplicated by perceptual hash.  Location is reported generically as
    "Document body" because DOCX doesn't expose per-paragraph image positions
    through a simple API.

    Returns the same dict schema as extract_images_pptx.
    """
    seen: dict[str, dict] = {}
    with zipfile.ZipFile(BytesIO(docx_bytes)) as z:
        for name in z.namelist():
            if name.lower().startswith("word/media/"):
                try:
                    data = z.read(name)
                    key  = str(get_image_hash(data))
                    if key not in seen:
                        seen[key] = {
                            "bytes":    data,
                            "zip_path": name,
                            "locations": ["Document body"],
                            "hash_str": key,
                        }
                except Exception:
                    pass
    return list(seen.values())


def extract_images_pdf(pdf_bytes: bytes) -> list[dict]:
    """
    Extract all unique embedded images from a PDF using PyMuPDF (fitz).

    PDFs store images as XObjects referenced by an integer `xref`.  We iterate
    every page's image list and de-duplicate by perceptual hash, accumulating
    all page numbers where each image appears.  The `xref` is stored in the
    dict so `embed_image_in_document` can call `doc.replace_image(xref, ...)`.

    Returns the same dict schema as extract_images_pptx (using `xref` instead
    of `zip_path`).
    """
    import fitz
    seen: dict[str, dict] = {}
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        for page_num, page in enumerate(doc):
            for img_info in page.get_images(full=True):
                xref = img_info[0]   # unique integer identifier for this image XObject
                try:
                    data = doc.extract_image(xref)["image"]
                    key  = str(get_image_hash(data))
                    loc  = f"Page {page_num + 1}"
                    if key not in seen:
                        seen[key] = {
                            "bytes":    data,
                            "xref":     xref,
                            "locations": [loc],
                            "hash_str": key,
                        }
                    elif loc not in seen[key]["locations"]:
                        seen[key]["locations"].append(loc)
                except Exception:
                    pass
    finally:
        doc.close()
    return list(seen.values())


def extract_images_any(file_bytes: bytes, filename: str) -> list[dict]:
    """
    Route image extraction to the correct handler based on file extension.
    Raises ValueError for unsupported formats.
    """
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".pptx":
        return extract_images_pptx(file_bytes)
    if ext in (".docx", ".doc"):
        return extract_images_docx(file_bytes)
    if ext == ".pdf":
        return extract_images_pdf(file_bytes)
    raise ValueError(f"Unsupported: {ext}")


# ============================================================
# RE-EMBED A (modified) IMAGE BACK INTO THE DOCUMENT
# ============================================================

def _zip_replace(zip_bytes: bytes, replacements: dict[str, bytes]) -> bytes:
    """
    Rebuild a ZIP archive, substituting specific entries with new byte content.

    This is the core primitive for safe PPTX/DOCX editing: rather than parsing
    and re-serialising XML, we treat the Office file as a ZIP and swap only the
    named entries.  Everything else — XML, relationships, themes, fonts — is
    copied byte-for-byte, guaranteeing no accidental modifications.

    Parameters
    ----------
    zip_bytes    : The original ZIP file bytes.
    replacements : Mapping of {internal_zip_path: new_bytes}.
    """
    buf = BytesIO()
    with zipfile.ZipFile(BytesIO(zip_bytes)) as zin, \
         zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            if item.filename in replacements:
                # Write the caller-supplied replacement bytes
                zout.writestr(item.filename, replacements[item.filename])
            else:
                # Copy unchanged entry verbatim
                zout.writestr(item.filename, zin.read(item.filename))
    return buf.getvalue()


def embed_image_in_document(
    file_bytes: bytes,
    filename: str,
    img_meta: dict,
    new_img_bytes: bytes,
) -> bytes:
    """
    Re-embed a (possibly region-modified) image back into its parent document.

    For PPTX/DOCX: delegates to `_zip_replace` using `img_meta["zip_path"]`.
    For PDF:       uses PyMuPDF's `replace_image(xref, ...)` API.

    Called by the region-detection flow after `replace_region_in_image` has
    painted the new logo over the detected bounding box.
    """
    ext = os.path.splitext(filename.lower())[1]

    if ext in (".pptx", ".docx", ".doc"):
        # Swap only the specific media file; all XML left untouched
        return _zip_replace(file_bytes, {img_meta["zip_path"]: new_img_bytes})

    if ext == ".pdf":
        import fitz
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        try:
            doc.replace_image(img_meta["xref"], stream=new_img_bytes)
            # garbage=4 removes all orphaned objects; deflate=True recompresses
            return doc.tobytes(garbage=4, deflate=True)
        finally:
            doc.close()

    raise ValueError(f"Unsupported: {ext}")


# ============================================================
# STANDARD HASH-BASED FULL-IMAGE REPLACEMENT (existing flow)
# ============================================================

def _apply_branding_rename_zip(zip_bytes: bytes) -> bytes:
    """
    Walk every XML file in a PPTX/DOCX ZIP archive and rename any
    capitalisation of "LTIMindtree" to "LTM".  Handles split runs the same
    way as `_replace_text_in_xml` (called with no old/new_text so only the
    always-on LTIMindtree pass fires).  Non-XML entries are copied unchanged.
    """
    buf = BytesIO()
    with zipfile.ZipFile(BytesIO(zip_bytes)) as zin, \
         zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.lower().endswith(".xml"):
                modified, n = _replace_text_in_xml(data)   # LTIMindtree-only pass
                if n:
                    data = modified
            zout.writestr(item.filename, data)
    return buf.getvalue()


def _apply_branding_rename_pdf(pdf_bytes: bytes) -> bytes:
    """
    Search every page of a PDF for visible text matching "LTIMindtree" in any
    capitalisation and redact it, inserting "LTM" in its place.

    Uses PyMuPDF's redaction annotation API:  the matched rectangle is whited
    out and the replacement string is drawn at the same position.  Font and
    size may not exactly match the original — this is an inherent limitation
    of PDF text editing without embedded font metrics.
    """
    import fitz
    # Enumerate the casing variants that the user explicitly mentioned plus
    # all-caps and sentence-case to cover common possibilities.
    _VARIANTS = ["LTIMindtree", "LTIMINDTREE", "ltimindtree", "Ltimindtree"]
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        for page in doc:
            found_any = False
            for variant in _VARIANTS:
                for rect in page.search_for(variant):
                    page.add_redact_annot(rect, text="LTM")
                    found_any = True
            if found_any:
                page.apply_redactions()
        return doc.tobytes(garbage=4, deflate=True)
    finally:
        doc.close()


def replace_logo_pptx(pptx_bytes, old_logo_bytes, new_logo_bytes, threshold):
    """
    Replace matching logos in a PPTX by swapping image bytes directly inside
    the ZIP (ppt/media/*).  The slide/master/layout XML is never touched, so
    all shapes, graphics, z-order, animations and formatting stay intact.

    How matching works:
      Each media file's perceptual hash is compared to the old logo's hash.
      If the Hamming distance is ≤ threshold, the bytes are replaced with the
      new logo (converted to PNG for consistency).

    Why not use python-pptx to remove/re-add the shape?
      Removing and re-inserting a shape via python-pptx changes the XML
      structure, resets z-order, and can corrupt master/layout relationships.
      Direct ZIP byte-swapping avoids all of that.
    """
    old_hash = get_image_hash(old_logo_bytes)
    new_png  = shrink_logo(new_logo_bytes)   # resize to 54% then normalise to RGBA PNG
    buf, count = BytesIO(), 0
    with zipfile.ZipFile(BytesIO(pptx_bytes)) as zin, \
         zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.lower().startswith("ppt/media/"):
                try:
                    if old_hash - get_image_hash(data) <= threshold:
                        # Hash distance within tolerance → this is the target logo
                        zout.writestr(item.filename, new_png)
                        count += 1
                        continue
                except Exception:
                    pass
            # All other entries (XML, relationships, themes…) copied unchanged
            zout.writestr(item.filename, data)
    # Also rename LTIMindtree → LTM in all slide/master XML text
    return _apply_branding_rename_zip(buf.getvalue()), count


def replace_logo_docx(docx_bytes, old_logo_bytes, new_logo_bytes, threshold):
    """
    Replace matching logos in a DOCX by swapping bytes inside the ZIP
    (`word/media/*`).  Same hash-based strategy as replace_logo_pptx.
    """
    old_hash = get_image_hash(old_logo_bytes)
    new_png  = shrink_logo(new_logo_bytes)   # resize to 54% then normalise to RGBA PNG
    buf, count = BytesIO(), 0
    with zipfile.ZipFile(BytesIO(docx_bytes)) as zin, \
         zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.lower().startswith("word/media/"):
                try:
                    if old_hash - get_image_hash(data) <= threshold:
                        zout.writestr(item.filename, new_png)
                        count += 1
                        continue
                except Exception:
                    pass
            zout.writestr(item.filename, data)
    # Also rename LTIMindtree → LTM in all XML text
    return _apply_branding_rename_zip(buf.getvalue()), count


def _to_pdf_safe_png(img_bytes: bytes) -> bytes:
    """
    Convert image bytes to an RGB PNG with no alpha channel, safe for PDFs.

    PDFs represent transparency via a separate /SMask (soft-mask) XObject.
    PyMuPDF's `replace_image` does not auto-create an SMask when given an RGBA
    PNG, causing the 4th (alpha) byte to be misinterpreted as a colour channel
    and producing a corrupted image in the output PDF.

    Fix: alpha-composite the new logo onto an opaque white background and
    save as a standard RGB PNG.  Logos that were already opaque are unaffected;
    logos with transparency get a clean white fill, which is the right default
    for document backgrounds.
    """
    img = Image.open(BytesIO(img_bytes)).convert("RGBA")
    bg  = Image.new("RGB", img.size, (255, 255, 255))
    bg.paste(img, mask=img.split()[3])   # alpha-composite RGBA onto white RGB
    buf = BytesIO()
    bg.save(buf, format="PNG")
    return buf.getvalue()


def replace_logo_pdf(pdf_bytes, old_logo_bytes, new_logo_bytes, threshold):
    """
    Replace matching logos in a PDF using PyMuPDF's replace_image API.

    Iterates every image XObject across all pages (de-duped by xref so shared
    images are only replaced once).  Matching is done by perceptual hash.
    The replacement image is converted to a PDF-safe RGB PNG via
    `_to_pdf_safe_png` to avoid alpha-channel corruption.

    `doc.tobytes(garbage=4, deflate=True)` rebuilds the PDF with all orphaned
    objects removed (garbage=4) and streams recompressed (deflate=True).
    """
    import fitz
    old_hash    = get_image_hash(old_logo_bytes)
    new_pdf_png = _to_pdf_safe_png(shrink_logo(new_logo_bytes))   # shrink then convert to RGB PNG
    count, done = 0, set()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        for page in doc:
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                if xref in done:
                    continue   # already replaced in a previous page reference
                try:
                    data = doc.extract_image(xref)["image"]
                    if old_hash - get_image_hash(data) <= threshold:
                        doc.replace_image(xref, stream=new_pdf_png)
                        done.add(xref)
                        count += 1
                except Exception:
                    pass
        result = doc.tobytes(garbage=4, deflate=True)
        return _apply_branding_rename_pdf(result), count
    finally:
        doc.close()


def replace_logo_any(file_bytes, filename, old_logo_bytes, new_logo_bytes, threshold):
    """Dispatch logo replacement to the correct format handler."""
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".pptx":
        return replace_logo_pptx(file_bytes, old_logo_bytes, new_logo_bytes, threshold)
    if ext in (".docx", ".doc"):
        return replace_logo_docx(file_bytes, old_logo_bytes, new_logo_bytes, threshold)
    if ext == ".pdf":
        return replace_logo_pdf(file_bytes, old_logo_bytes, new_logo_bytes, threshold)
    raise ValueError(f"Unsupported: {ext}")


def replace_multiple_logos_any(
    file_bytes: bytes,
    filename: str,
    old_logos_bytes: list[bytes],   # one entry per selected logo
    new_logo_bytes: bytes,
    threshold: int = 8,
) -> tuple[bytes, int]:
    """
    Replace every logo in `old_logos_bytes` inside a single document in one pass.

    Applies `replace_logo_any` sequentially for each reference logo image,
    feeding the output of one replacement as the input to the next.  This
    allows multiple distinct logos (e.g. header + footer logo) to all be
    replaced in a single call that returns one combined output file.

    Returns (updated_file_bytes, total_replacement_count).
    """
    data  = file_bytes
    total = 0
    for old_bytes in old_logos_bytes:
        data, n = replace_logo_any(data, filename, old_bytes, new_logo_bytes, threshold)
        total += n
    return data, total


# ============================================================
# TEXT REPLACEMENT  (handles split runs across a:r / w:r)
# ============================================================

def _replace_text_in_xml(
    xml_bytes: bytes,
    old_text: str,
    new_text: str,
    extra_pairs: list[tuple[str, str]] | None = None,
) -> tuple[bytes, int]:
    """
    Find and replace text within a single PPTX/DOCX XML part, correctly
    handling cases where PowerPoint/Word has split the target string across
    multiple text run elements (`<a:r>` / `<w:r>`).

    Parameters
    ----------
    xml_bytes    : Raw bytes of the XML part.
    old_text     : Primary string to search for.
    new_text     : Replacement for old_text.
    extra_pairs  : Additional (old, new) string pairs applied after the primary
                   replacement (e.g. for year normalisation).  Each pair is
                   applied as a plain string replacement on the already-combined
                   paragraph text.

    How split runs are handled:
      Office stores visible strings across multiple <a:r>/<w:r> run elements.
      Each paragraph's run texts are concatenated, the replacement is applied
      to the combined string, the result is written into the first run's text
      node, and remaining run text nodes are cleared.  All formatting properties
      (<a:rPr>/<w:rPr>) are preserved.

    Returns (modified_xml_bytes, replacement_count).
    If no changes are made the original bytes are returned (count=0).
    """
    import re
    from lxml import etree
    # DrawingML namespace (PPTX text)
    _A     = "http://schemas.openxmlformats.org/drawingml/2006/main"
    # WordprocessingML namespace (DOCX text)
    _W     = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    # xml:space attribute used to preserve leading/trailing whitespace
    _SPACE = "{http://www.w3.org/XML/1998/namespace}space"
    count  = 0

    try:
        root = etree.fromstring(xml_bytes)
    except Exception:
        return xml_bytes, 0   # unparseable XML — return untouched

    # Process both DrawingML (PPTX) and WordprocessingML (DOCX) paragraph types
    for p_tag, r_tag, t_tag in [
        (f"{{{_A}}}p", f"{{{_A}}}r", f"{{{_A}}}t"),   # PPTX
        (f"{{{_W}}}p", f"{{{_W}}}r", f"{{{_W}}}t"),   # DOCX
    ]:
        for para in root.iter(p_tag):
            runs   = para.findall(r_tag)
            # Collect only those runs that have a text element
            telems = [r.find(t_tag) for r in runs]
            telems = [t for t in telems if t is not None]
            if not telems:
                continue
            # Concatenate all run texts to get the full visible string
            full = "".join(t.text or "" for t in telems)
            original_full = full

            # Apply the primary replacement
            if old_text in full:
                full   = full.replace(old_text, new_text)
                count += original_full.count(old_text)

            # Apply any extra (old, new) pairs — e.g. year normalisation
            if extra_pairs:
                for ep_old, ep_new in extra_pairs:
                    if ep_old in full:
                        full   = full.replace(ep_old, ep_new)
                        count += 1

            if full == original_full:
                continue   # nothing changed in this paragraph

            # Write entire result into first run's text node
            telems[0].text = full
            # Preserve leading/trailing spaces if present
            if full != full.strip():
                telems[0].set(_SPACE, "preserve")
            # Clear the remaining fragment text nodes
            for t in telems[1:]:
                t.text = ""

    if count == 0:
        return xml_bytes, 0
    # Re-serialise with XML declaration to maintain byte-level compatibility
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True), count


def _build_year_pairs(new_text: str, target_year: str = "2026") -> list[tuple[str, str]]:
    """
    Build extra (old, new) replacement pairs so that any 4-digit year that
    follows a known suffix in `new_text` is updated to `target_year`.

    For example if new_text ends with "LTM | Privileged and Confidential",
    this returns pairs like:
        ("LTM | Privileged and Confidential 2021", "LTM | Privileged and Confidential 2026")
        ("LTM | Privileged and Confidential 2022", "LTM | Privileged and Confidential 2026")
        ... up to the year before target_year

    We also cover the *old* branding string so the year gets updated regardless
    of whether the primary text replacement fires first or not in the same para.
    """
    import re
    pairs: list[tuple[str, str]] = []
    suffixes = [new_text]  # catch year on already-replaced new text
    for year in range(2018, int(target_year)):   # 2018–2025 covers realistic range
        yr = str(year)
        for sfx in suffixes:
            candidate_old = f"{sfx} {yr}"
            candidate_new = f"{sfx} {target_year}"
            if candidate_old != candidate_new:
                pairs.append((candidate_old, candidate_new))
    return pairs


def replace_text_pptx(
    pptx_bytes: bytes,
    old_text: str,
    new_text: str,
    extra_pairs: list[tuple[str, str]] | None = None,
) -> tuple[bytes, int]:
    """
    Replace text across all XML parts of a PPTX that can contain visible text:
    slides, slide layouts, slide masters, notes slides, handout masters, and
    notes masters.

    Non-XML entries (media files, theme files, font embeds, rels) are copied
    byte-for-byte.  Returns (updated_pptx_bytes, total_replacement_count).
    """
    # All ZIP directories that contain XML with visible text content
    XML_DIRS = (
        "ppt/slides/", "ppt/slidelayouts/", "ppt/slidemasters/",
        "ppt/notesslides/", "ppt/handoutmasters/", "ppt/notesmasters/",
    )
    total = 0
    buf   = BytesIO()
    with zipfile.ZipFile(BytesIO(pptx_bytes)) as zin, \
         zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data  = zin.read(item.filename)
            fname = item.filename.lower()
            if fname.endswith(".xml") and any(fname.startswith(d) for d in XML_DIRS):
                modified, n = _replace_text_in_xml(data, old_text, new_text, extra_pairs)
                if n:
                    data   = modified
                    total += n
            zout.writestr(item.filename, data)
    return buf.getvalue(), total


def replace_text_docx(
    docx_bytes: bytes,
    old_text: str,
    new_text: str,
    extra_pairs: list[tuple[str, str]] | None = None,
) -> tuple[bytes, int]:
    """
    Replace text across all XML parts in a DOCX.

    Processes every `.xml` file in the archive (document body, headers,
    footers, footnotes, endnotes, etc.).  Non-XML entries are copied unchanged.
    Returns (updated_docx_bytes, total_replacement_count).
    """
    total = 0
    buf   = BytesIO()
    with zipfile.ZipFile(BytesIO(docx_bytes)) as zin, \
         zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.lower().endswith(".xml"):
                modified, n = _replace_text_in_xml(data, old_text, new_text, extra_pairs)
                if n:
                    data   = modified
                    total += n
            zout.writestr(item.filename, data)
    return buf.getvalue(), total


def replace_text_any(
    file_bytes: bytes,
    filename: str,
    old_text: str,
    new_text: str,
    update_year: bool = True,
    target_year: str = "2026",
) -> tuple[bytes, int]:
    """
    Dispatch text replacement to the correct format handler.

    When `update_year=True` (default), also replaces any occurrence of
    `new_text + " YYYY"` (for years 2018–2025) with `new_text + " 2026"`
    so that the copyright/confidentiality year is normalised in the same pass.
    """
    extra_pairs = _build_year_pairs(new_text, target_year) if update_year else None
    ext = os.path.splitext(filename.lower())[1]
    if ext == ".pptx":
        return replace_text_pptx(file_bytes, old_text, new_text, extra_pairs)
    if ext in (".docx", ".doc"):
        return replace_text_docx(file_bytes, old_text, new_text, extra_pairs)
    raise ValueError(f"Text replacement not supported for {ext}")


def process_zip(zip_bytes, old_logo_bytes, new_logo_bytes, threshold):
    """
    Apply hash-based whole-image logo replacement to every supported document
    in a ZIP archive in a single pass.

    Iterates all entries; for PPTX/DOCX/PDF files calls `replace_logo_any`,
    all other files (images, fonts, other ZIPs) are copied unchanged.

    Returns (output_zip_bytes, total_replacements, files_processed).
    """
    buf, total, files = BytesIO(), 0, 0
    with zipfile.ZipFile(BytesIO(zip_bytes)) as zin, \
         zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            ext  = os.path.splitext(item.filename.lower())[1]
            if ext in (".pptx", ".docx", ".doc", ".pdf"):
                try:
                    updated, n = replace_logo_any(data, item.filename, old_logo_bytes, new_logo_bytes, threshold)
                    zout.writestr(item.filename, updated)
                    total += n
                    files += 1
                    continue
                except Exception:
                    pass
            zout.writestr(item.filename, data)
    return buf.getvalue(), total, files


def process_zip_region(
    zip_outer_bytes: bytes,
    detect_fn,          # callable: (img_bytes: bytes) -> dict|None  {x1,y1,x2,y2}
    new_logo_bytes: bytes,
    progress_cb=None,   # optional callable(msg: str) for live status updates
) -> tuple[bytes, int, int, int]:
    """
    Automated region-based logo replacement across every supported document in a ZIP.

    This handles the hard case where a logo is *baked into* a background image
    rather than existing as a standalone picture shape.  For each embedded image
    in each document, a detection function (OCR or template match) is run to
    locate the logo bbox, then `replace_region_in_image` paints the new logo
    over that region, and `embed_image_in_document` writes it back.

    Parameters
    ----------
    zip_outer_bytes : Raw bytes of the outer ZIP file.
    detect_fn       : Callable receiving image bytes, returning a bbox dict
                      {x1, y1, x2, y2} or None if not found.
    new_logo_bytes  : The new logo to paint into each detected region.
    progress_cb     : Optional callback for real-time progress messages in the UI.

    Returns
    -------
    (output_zip_bytes, files_processed, regions_replaced, files_with_no_match)
    """
    DOC_EXTS = {".pptx", ".docx", ".doc", ".pdf"}
    out_buf          = BytesIO()
    files_processed  = 0
    regions_replaced = 0
    files_no_match   = 0

    with zipfile.ZipFile(BytesIO(zip_outer_bytes)) as zin, \
         zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:

        all_items = zin.infolist()
        for idx, item in enumerate(all_items):
            data = zin.read(item.filename)
            ext  = os.path.splitext(item.filename.lower())[1]

            if ext not in DOC_EXTS:
                zout.writestr(item.filename, data)
                continue

            if progress_cb:
                short = os.path.basename(item.filename)
                progress_cb(f"[{idx+1}/{len(all_items)}] Processing {short}…")

            try:
                images    = extract_images_any(data, item.filename)
                file_data = data
                found_in_file = 0

                for img_meta in images:
                    try:
                        bbox = detect_fn(img_meta["bytes"])
                    except Exception:
                        bbox = None

                    if bbox is None:
                        continue   # logo not found in this image

                    # Paint new logo over the detected bounding box
                    modified_img = replace_region_in_image(
                        img_meta["bytes"],
                        bbox["x1"], bbox["y1"], bbox["x2"], bbox["y2"],
                        new_logo_bytes,
                    )
                    img_meta["bytes"] = modified_img   # update meta for chained detection
                    file_data = embed_image_in_document(
                        file_data, item.filename, img_meta, modified_img
                    )
                    found_in_file    += 1
                    regions_replaced += 1

                if found_in_file == 0:
                    files_no_match += 1
                    if progress_cb:
                        progress_cb(f"  ↳ No logo found in {os.path.basename(item.filename)}")
                else:
                    if progress_cb:
                        progress_cb(
                            f"  ↳ Replaced {found_in_file} region(s) "
                            f"in {os.path.basename(item.filename)}"
                        )

                zout.writestr(item.filename, file_data)
                files_processed += 1

            except Exception as e:
                if progress_cb:
                    progress_cb(f"  ↳ ERROR on {os.path.basename(item.filename)}: {e}")
                zout.writestr(item.filename, data)   # pass through unchanged on error

    return out_buf.getvalue(), files_processed, regions_replaced, files_no_match


def process_zip_by_refs(
    zip_bytes: bytes,
    ref_images: list[bytes],   # reference logo images selected from the sample scan
    new_logo_bytes: bytes,
    threshold: int = 8,        # perceptual hash distance tolerance
    progress_cb=None,
) -> tuple[bytes, int, int]:
    """
    Replace logos across all documents in a ZIP using pre-selected reference images.

    This is the "Guided" ZIP flow.  The user first scans one sample file to see
    its embedded images, selects which ones are logos, and those image bytes
    become `ref_images`.  Every document in the ZIP is then processed: any
    embedded image whose pHash is within `threshold` of any reference hash is
    replaced with the new logo.

    Parameters
    ----------
    zip_bytes   : Raw bytes of the ZIP archive.
    ref_images  : List of raw bytes for each reference logo the user selected.
    new_logo_bytes : The new logo bytes to embed.
    threshold   : Maximum pHash Hamming distance to count as a match (0–20).
                  0 = exact match only.  8 = tolerates minor compression changes.
    progress_cb : Optional callback for live log output in the UI.

    Returns
    -------
    (output_zip_bytes, files_processed, total_replacements)
    """
    # Pre-compute hashes of all reference logos once for efficiency
    ref_hashes = [get_image_hash(b) for b in ref_images]
    out_buf      = BytesIO()
    files_proc   = 0
    total_repl   = 0
    DOC_EXTS     = {".pptx", ".docx", ".doc", ".pdf"}

    with zipfile.ZipFile(BytesIO(zip_bytes)) as zin, \
         zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:

        all_items = zin.infolist()
        for idx, item in enumerate(all_items):
            data = zin.read(item.filename)
            ext  = os.path.splitext(item.filename.lower())[1]

            if ext not in DOC_EXTS:
                zout.writestr(item.filename, data)
                continue

            short = os.path.basename(item.filename)
            if progress_cb:
                progress_cb(f"[{idx+1}/{len(all_items)}] Processing {short}…")

            try:
                images    = extract_images_any(data, item.filename)
                file_data = data
                count     = 0

                for img_meta in images:
                    try:
                        h = get_image_hash(img_meta["bytes"])
                    except Exception:
                        continue

                    # Check if this image matches any user-selected reference logo
                    if any(abs(h - rh) <= threshold for rh in ref_hashes):
                        file_data = embed_image_in_document(
                            file_data, item.filename, img_meta,
                            to_png(new_logo_bytes),
                        )
                        count      += 1
                        total_repl += 1

                if progress_cb:
                    if count:
                        progress_cb(f"  ↳ Replaced {count} logo image(s)")
                    else:
                        progress_cb(f"  ↳ No matching logos found")

                zout.writestr(item.filename, file_data)
                files_proc += 1

            except Exception as e:
                if progress_cb:
                    progress_cb(f"  ↳ ERROR: {e}")
                zout.writestr(item.filename, data)

    return out_buf.getvalue(), files_proc, total_repl


# ============================================================
# SESSION STATE
# ============================================================
# Streamlit reruns the entire script on every user interaction.
# `st.session_state` is a dict-like object that persists across reruns,
# allowing us to remember which file is loaded, which images were scanned,
# and which logos the user has selected.
# ============================================================

_defaults = {
    "images":          [],      # list of image-meta dicts from the last scan
    "selected":        None,    # index of the image currently open in the Inspect panel
    "doc_cache":       None,    # tuple (filename, bytes) — avoids re-reading on every rerun
    "bbox":            None,    # detected region dict {x1,y1,x2,y2} for intra-image mode
    "drill_active":    False,   # whether the "Find inside image" drill-down panel is open
    # ZIP guided flow
    "zip_sample_imgs": [],      # image-meta dicts scanned from the chosen sample file
    "zip_ref_indices": [],      # list of indices into zip_sample_imgs selected by the user
    # single-file multi-select
    "selected_set":    set(),   # set of indices of images marked for bulk replacement
}
# Initialise any missing keys without overwriting values already in state
for _k, _v in _defaults.items():
    if _k not in st.session_state:
        st.session_state[_k] = _v


# ============================================================
# UI HELPERS
# ============================================================

# Maps file extensions to their official MIME types for the browser download
MIME = {
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".doc":  "application/msword",
    ".pdf":  "application/pdf",
}


def _download_doc(label: str, data: bytes, stem: str, ext: str):
    """
    Render a Streamlit download button for an updated document.
    The output filename appends `_updated` to the original stem to avoid
    overwriting the source file accidentally.
    """
    st.download_button(
        label, data,
        file_name=f"{stem}_updated{ext}",
        mime=MIME.get(ext, "application/octet-stream"),
    )


# ============================================================
# PAGE
# ============================================================

st.set_page_config(page_title="Document Logo Replacer", page_icon="🖼️", layout="wide")
st.title("🖼️ Document Logo Replacer")
st.caption(
    "Supports PPTX · DOCX · PDF · ZIP batch. "
    "Finds logos on slide masters and layouts. "
    "Can also detect logos **baked into** a template background image using OCR or template matching."
)
st.divider()

# ── Upload ────────────────────────────────────────────────────────────────────

doc_file = st.file_uploader(
    "Upload document",
    type=["pptx", "docx", "doc", "pdf", "zip"],
    help="Single file or a ZIP containing multiple documents.",
)
# Detect whether the uploaded file is a ZIP batch
is_zip = doc_file is not None and doc_file.name.lower().endswith(".zip")

if doc_file:
    cache = st.session_state.doc_cache
    # Only re-read file bytes when a *different* file is uploaded;
    # on subsequent reruns (e.g. button clicks) the cached bytes are reused
    if cache is None or cache[0] != doc_file.name:
        doc_file.seek(0)
        st.session_state.doc_cache        = (doc_file.name, doc_file.read())
        # Reset all scan/selection state for the new file
        st.session_state.images           = []
        st.session_state.selected         = None
        st.session_state.selected_set     = set()
        st.session_state.bbox             = None
        st.session_state.drill_active     = False
        st.session_state.zip_sample_imgs  = []
        st.session_state.zip_ref_indices  = []

# ── ZIP flow ──────────────────────────────────────────────────────────────────

if is_zip:
    name_zip, zbytes_zip = st.session_state.doc_cache
    stem_zip = os.path.splitext(name_zip)[0]

    # Count supported files in the ZIP for display
    with zipfile.ZipFile(BytesIO(zbytes_zip)) as _z:
        _supported = [
            i.filename for i in _z.infolist()
            if os.path.splitext(i.filename.lower())[1] in (".pptx",".docx",".doc",".pdf")
        ]
    st.info(
        f"ZIP contains **{len(_supported)}** supported file(s): "
        + ", ".join(os.path.basename(f) for f in _supported[:8])
        + (" …" if len(_supported) > 8 else "")
    )

    zip_tab_hash, zip_tab_guided, zip_tab_region = st.tabs([
        "🔄 Quick — upload old logo file",
        "🎯 Guided — scan sample & select logos",
        "🔬 Region detect — logos baked into backgrounds",
    ])

    # ── ZIP Tab A: hash-based ─────────────────────────────────────────────────
    with zip_tab_hash:
        st.markdown(
            "Upload the **old logo** as a reference image. "
            "Every embedded image across all files that matches it will be swapped out wholesale."
        )
        zh_c1, zh_c2 = st.columns(2)
        with zh_c1:
            old_logo_hash = st.file_uploader(
                "Old logo (reference)", type=["png","jpg","jpeg","bmp","webp"], key="zip_hash_old"
            )
            if old_logo_hash:
                st.image(old_logo_hash, use_container_width=True)
        with zh_c2:
            new_logo_hash = st.file_uploader(
                "New logo", type=["png","jpg","jpeg","bmp","webp"], key="zip_hash_new"
            )
            if new_logo_hash:
                st.image(new_logo_hash, use_container_width=True)

        thr_hash = st.slider("Match sensitivity", 0, 20, 5, key="zip_hash_thresh")

        if st.button("Replace Logos (whole-image) in all files", type="primary",
                     disabled=not (old_logo_hash and new_logo_hash), key="btn_zip_hash"):
            with st.spinner("Processing ZIP…"):
                out_zip, total, files = process_zip(
                    zbytes_zip, old_logo_hash.read(), new_logo_hash.read(), thr_hash
                )
            if files == 0:
                st.warning("No supported files found in the ZIP.")
            elif total == 0:
                st.warning(f"Processed {files} file(s) but no matching logos found. Try raising sensitivity.")
            else:
                st.success(f"Replaced **{total}** logo(s) across **{files}** file(s).")
            st.download_button(
                "⬇️ Download updated ZIP", out_zip,
                file_name=stem_zip + "_updated.zip", mime="application/zip"
            )

    # ── ZIP Tab B: guided scan + multi-select ────────────────────────────────
    with zip_tab_guided:
        st.markdown(
            "Scan one file from the ZIP to see all embedded images, "
            "then **select which images are logos** you want replaced. "
            "The tool will find those logos in **every file in the ZIP** and swap them "
            "for the new logo automatically. You can select multiple logos "
            "(e.g. a title-slide logo **and** a footer logo)."
        )
        st.divider()

        # Step 1 ──────────────────────────────────────────────────────────────
        st.markdown("#### Step 1 — Scan a sample file from the ZIP")
        _sample_names  = [os.path.basename(f) for f in _supported]
        _sample_choice = st.selectbox(
            "Pick which file to use as a reference scan",
            _sample_names, key="zip_sample_choice",
        )

        if st.button("🔍 Scan sample file", key="btn_zip_scan"):
            with zipfile.ZipFile(BytesIO(zbytes_zip)) as _zs:
                _full_paths = [f for f in _supported if os.path.basename(f) == _sample_choice]
                if _full_paths:
                    _raw = _zs.read(_full_paths[0])
                    with st.spinner(f"Scanning {_sample_choice}…"):
                        _scanned = extract_images_any(_raw, _full_paths[0])
                    st.session_state.zip_sample_imgs  = _scanned
                    st.session_state.zip_ref_indices  = []
                    if not _scanned:
                        st.warning("No images found in this file. Try picking a different sample.")
                    st.rerun()

        # Step 2 ──────────────────────────────────────────────────────────────
        _sample_imgs = st.session_state.zip_sample_imgs
        if _sample_imgs:
            st.divider()
            st.markdown(
                "#### Step 2 — Select the logo(s) to replace\n"
                "Click **Select as logo** to toggle an image. "
                "Select as many as needed — e.g. one header logo and one footer logo."
            )
            _COLS = 5
            for _ri, _row in enumerate(
                [_sample_imgs[i : i + _COLS] for i in range(0, len(_sample_imgs), _COLS)]
            ):
                _cols = st.columns(_COLS)
                for _ci, (_col, _img) in enumerate(zip(_cols, _row)):
                    _idx    = _ri * _COLS + _ci
                    _is_sel = _idx in st.session_state.zip_ref_indices
                    _border = "2px solid #00c853" if _is_sel else "2px solid #ddd"
                    with _col:
                        st.markdown(
                            f'<div style="border:{_border};border-radius:6px;padding:4px">',
                            unsafe_allow_html=True,
                        )
                        st.image(thumbnail_bytes(_img["bytes"]), use_container_width=True)
                        _locs  = _img.get("locations", [])
                        _label = (
                            _locs[0] if len(_locs) == 1
                            else (f"{_locs[0]} +{len(_locs)-1} more" if _locs else "—")
                        )
                        st.caption(_label)
                        _btn_lbl = "✅ Selected" if _is_sel else "Select as logo"
                        if st.button(_btn_lbl, key=f"zg_{_idx}", use_container_width=True):
                            _refs = list(st.session_state.zip_ref_indices)
                            if _is_sel:
                                _refs.remove(_idx)
                            else:
                                _refs.append(_idx)
                            st.session_state.zip_ref_indices = _refs
                            st.rerun()
                        st.markdown("</div>", unsafe_allow_html=True)

            _n_sel = len(st.session_state.zip_ref_indices)
            if _n_sel:
                st.success(f"{_n_sel} logo image(s) selected for replacement.")
            else:
                st.info("Click images above to mark them as logos to be replaced.")

        # Step 3 ──────────────────────────────────────────────────────────────
        _refs_idx = st.session_state.zip_ref_indices
        if _refs_idx and _sample_imgs:
            st.divider()
            st.markdown("#### Step 3 — Upload new logo & run")
            _gc1, _gc2 = st.columns(2)
            with _gc1:
                st.markdown("**Logo(s) that will be replaced (your selection):**")
                _sel_thumbs = [_sample_imgs[i] for i in _refs_idx]
                _tcols      = st.columns(min(len(_sel_thumbs), 4))
                for _tc, _si in zip(_tcols, _sel_thumbs):
                    with _tc:
                        st.image(thumbnail_bytes(_si["bytes"]), use_container_width=True)
            with _gc2:
                st.markdown("**New logo to apply to all selected logos:**")
                zg_new_logo = st.file_uploader(
                    "New logo",
                    type=["png", "jpg", "jpeg", "bmp", "webp"],
                    key="zip_guided_new",
                    label_visibility="collapsed",
                )
                if zg_new_logo:
                    st.image(zg_new_logo, use_container_width=True)

            zg_thresh = st.slider(
                "Match sensitivity — hash distance (higher = more tolerant of minor variations between files)",
                0, 20, 8, key="zip_guided_thresh",
            )

            if st.button(
                "🚀 Replace logos across all files in ZIP",
                type="primary",
                disabled=zg_new_logo is None,
                key="btn_zip_guided",
            ):
                _ref_bytes = [_sample_imgs[i]["bytes"] for i in _refs_idx]
                _new_bytes = zg_new_logo.read()

                _log_ph    = st.empty()
                _log_lines_g: list[str] = []

                def _glog(msg: str):
                    _log_lines_g.append(msg)
                    _log_ph.code("\n".join(_log_lines_g[-30:]))

                _glog("Starting guided batch replacement…")
                with st.spinner("Processing ZIP — this may take a moment…"):
                    _out_zip, _f_done, _r_done = process_zip_by_refs(
                        zbytes_zip, _ref_bytes, _new_bytes, zg_thresh, progress_cb=_glog
                    )
                _glog("")
                _glog(f"✅ Done — {_f_done} file(s) processed, {_r_done} replacement(s) made.")

                if _r_done == 0:
                    st.warning(
                        "No matching logos were found. "
                        "Try raising the sensitivity slider or re-selecting the logo images in Step 2."
                    )
                else:
                    st.success(
                        f"Replaced **{_r_done}** logo instance(s) across **{_f_done}** file(s)."
                    )
                st.download_button(
                    "⬇️ Download updated ZIP", _out_zip,
                    file_name=stem_zip + "_updated.zip", mime="application/zip",
                    key="dl_guided_zip",
                )

    # ── ZIP Tab C: automated region detection (baked-in logos) ───────────────
    with zip_tab_region:
        st.markdown(
            "Use this when the logo is **baked into a full-slide background image** "
            "and cannot be selected as a separate shape. "
            "The tool scans every image in every file and detects the logo region using "
            "OCR (text recognition) or visual template matching, then pastes the new logo over it."
        )
        st.divider()

        zr_method = st.radio(
            "Detection method",
            ["OCR — find by text", "Template match — find by visual crop"],
            horizontal=True, key="zip_region_method",
        )

        zr_c1, zr_c2 = st.columns(2)

        with zr_c1:
            st.markdown("**Detection settings**")
            if zr_method.startswith("OCR"):
                zr_ocr_text   = st.text_input("Text to search for", value="LTIMindtree", key="zr_ocr_text")
                zr_expand_l   = st.slider("Expand left (% of text width, to include icon)",
                                          0, 200, 90, key="zr_expand_l")
                zr_vpad       = st.slider("Vertical padding (%)", 0, 50, 15, key="zr_vpad")
                zr_tmpl_bytes = None
            else:
                st.markdown(
                    "📸 **A screenshot of the logo is fine.** "
                    "The tool will automatically remove any surrounding padding / desktop background "
                    "and normalise contrast before searching. "
                    "Try to capture the logo area as closely as possible for best results."
                )
                zr_tmpl_file  = st.file_uploader(
                    "Screenshot or crop of old logo",
                    type=["png","jpg","jpeg","bmp","webp"], key="zr_tmpl"
                )
                zr_tmpl_bytes = zr_tmpl_file.read() if zr_tmpl_file else None
                if zr_tmpl_bytes:
                    col_raw, col_crop = st.columns(2)
                    with col_raw:
                        st.caption("Your upload")
                        st.image(zr_tmpl_bytes, use_container_width=True)
                    with col_crop:
                        st.caption("After auto-crop (what will be matched)")
                        st.image(autocrop_screenshot(zr_tmpl_bytes), use_container_width=True)
                zr_min_score  = st.slider(
                    "Minimum match score (0=loose, 1=exact)", 0.0, 1.0, 0.40, step=0.05, key="zr_score"
                )

        with zr_c2:
            st.markdown("**New logo**")
            zr_new_logo = st.file_uploader(
                "New logo", type=["png","jpg","jpeg","bmp","webp"], key="zr_new_logo"
            )
            if zr_new_logo:
                st.image(zr_new_logo, use_container_width=True)

        # Readiness check
        zr_ocr_ready  = zr_method.startswith("OCR") and zr_new_logo is not None
        zr_tmpl_ready = zr_method.startswith("Template") and zr_tmpl_bytes is not None and zr_new_logo is not None
        zr_ready = zr_ocr_ready or zr_tmpl_ready

        if st.button(
            "🚀 Run automated replacement across all files in ZIP",
            type="primary", disabled=not zr_ready, key="btn_zip_region"
        ):
            zr_new_bytes = zr_new_logo.read()

            # Build the detect_fn from current settings
            if zr_method.startswith("OCR"):
                _text   = zr_ocr_text
                _expl   = zr_expand_l / 100
                _vpad   = zr_vpad / 100
                def _detect(img_bytes, _t=_text, _e=_expl, _v=_vpad):
                    return detect_by_ocr(img_bytes, _t, _e, _v)
            else:
                _tmpl   = zr_tmpl_bytes
                _score  = zr_min_score
                def _detect(img_bytes, _tmpl=_tmpl, _s=_score):
                    return detect_by_template(img_bytes, _tmpl, _s)

            log_placeholder = st.empty()
            log_lines: list[str] = []

            def _log(msg: str):
                log_lines.append(msg)
                log_placeholder.code("\n".join(log_lines[-30:]))  # show last 30 lines

            _log("Starting automated region replacement…")

            with st.spinner("Running detection across all files — this may take a few minutes…"):
                out_zip, files_done, regions_done, no_match = process_zip_region(
                    zbytes_zip, _detect, zr_new_bytes, progress_cb=_log
                )

            _log("")
            _log(f"✅ Done — {files_done} file(s) processed, "
                 f"{regions_done} region(s) replaced, "
                 f"{no_match} file(s) had no match.")

            if regions_done == 0:
                st.warning(
                    "No logo regions were detected in any file. "
                    "Try adjusting detection settings or switching methods."
                )
            else:
                st.success(
                    f"Replaced **{regions_done}** logo region(s) across "
                    f"**{files_done - no_match}** file(s)."
                )

            st.download_button(
                "⬇️ Download updated ZIP", out_zip,
                file_name=stem_zip + "_updated.zip", mime="application/zip"
            )

    st.stop()

# ── Scan ──────────────────────────────────────────────────────────────────────

if doc_file:
    if st.button("🔍 Scan document for all embedded images"):
        name, fbytes = st.session_state.doc_cache
        with st.spinner("Scanning…"):
            imgs = extract_images_any(fbytes, name)
        st.session_state.images     = imgs
        st.session_state.selected   = None
        st.session_state.selected_set = set()
        st.session_state.bbox       = None
        st.session_state.drill_active = False
        if not imgs:
            st.warning("No images found in this document.")

# ── Image grid ────────────────────────────────────────────────────────────────

if st.session_state.images:
    st.divider()
    n = len(st.session_state.images)
    _n_sel = len(st.session_state.selected_set)
    st.subheader(
        f"Step 2 \u2014 Mark logo(s) to replace  ({n} unique image{'s' if n != 1 else ''} found)"
    )
    st.caption(
        "**Mark as logo** to add an image to the replacement list (supports multi-select). "
        "**Inspect** to open the region-detect panel for a single image."
    )

    COLS = 5
    imgs = st.session_state.images
    for row_i, row in enumerate([imgs[i:i+COLS] for i in range(0, len(imgs), COLS)]):
        cols = st.columns(COLS)
        for col_i, (col, img) in enumerate(zip(cols, row)):
            idx        = row_i * COLS + col_i
            is_marked  = idx in st.session_state.selected_set
            is_inspect = st.session_state.selected == idx
            if is_marked:
                border = "2px solid #00c853"
            elif is_inspect:
                border = "2px solid #2196f3"
            else:
                border = "2px solid #ddd"
            with col:
                st.markdown(
                    f'<div style="border:{border};border-radius:6px;padding:4px">',
                    unsafe_allow_html=True,
                )
                st.image(thumbnail_bytes(img["bytes"]), use_container_width=True)
                locs  = img["locations"]
                label = locs[0] if len(locs) == 1 else f"{locs[0]} +{len(locs)-1} more"
                st.caption(label)
                bc1, bc2 = st.columns(2)
                with bc1:
                    mark_lbl = "\u2705 Marked" if is_marked else "Mark as logo"
                    if st.button(mark_lbl, key=f"mark_{idx}", use_container_width=True):
                        _s = set(st.session_state.selected_set)
                        if is_marked:
                            _s.discard(idx)
                        else:
                            _s.add(idx)
                        st.session_state.selected_set = _s
                        st.rerun()
                with bc2:
                    insp_lbl = "\U0001f535 Inspecting" if is_inspect else "\U0001f50d Inspect"
                    if st.button(insp_lbl, key=f"insp_{idx}", use_container_width=True):
                        st.session_state.selected     = idx
                        st.session_state.bbox         = None
                        st.session_state.drill_active = False
                        st.rerun()
                st.markdown("</div>", unsafe_allow_html=True)

    if _n_sel:
        st.success(f"{_n_sel} logo image{'s' if _n_sel != 1 else ''} marked for replacement.")
    else:
        st.info("Click **Mark as logo** on one or more images above, then upload the new logo below.")

# ── Step 3: Bulk replace all marked logos ───────────────────────────────────────

_marked_indices = list(st.session_state.selected_set)
if _marked_indices and st.session_state.images:
    _marked_imgs = [st.session_state.images[i] for i in sorted(_marked_indices)]
    name_sf, fbytes_sf = st.session_state.doc_cache
    ext_sf  = os.path.splitext(name_sf.lower())[1]
    stem_sf = os.path.splitext(name_sf)[0]

    st.divider()
    st.subheader("Step 3 \u2014 Replace selected logo(s)")

    _mc1, _mc2 = st.columns([2, 1])
    with _mc1:
        st.markdown("**Logo(s) that will be replaced:**")
        _tcols = st.columns(min(len(_marked_imgs), 5))
        for _tc, _mi in zip(_tcols, _marked_imgs):
            with _tc:
                st.image(thumbnail_bytes(_mi["bytes"]), use_container_width=True)
                _ml = _mi.get("locations", [])
                st.caption(_ml[0] if len(_ml) == 1 else (f"{_ml[0]} +{len(_ml)-1} more" if _ml else "—"))
    with _mc2:
        st.markdown("**New logo to apply:**")
        _new_logo_multi = st.file_uploader(
            "New logo",
            type=["png", "jpg", "jpeg", "bmp", "webp"],
            key="multi_new_logo",
            label_visibility="collapsed",
        )
        if _new_logo_multi:
            st.image(_new_logo_multi, use_container_width=True)

    _thr_multi = st.slider(
        "Match sensitivity (hash distance — higher = more tolerant)",
        0, 20, 8, key="multi_thresh",
    )

    with st.expander("\U0001f524 Also replace text in this document (optional)"):
        _tr_c1sf, _tr_c2sf = st.columns(2)
        with _tr_c1sf:
            _tr_old_sf = st.text_input(
                "Find text",
                value="LTIMindtree | Privileged and Confidential",
                key="tr_old_sf",
            )
        with _tr_c2sf:
            _tr_new_sf = st.text_input(
                "Replace with",
                value="LTM | Privileged and Confidential",
                key="tr_new_sf",
            )

    if st.button(
        f"\U0001f504 Replace {len(_marked_imgs)} logo image{'s' if len(_marked_imgs) != 1 else ''} in document",
        type="primary",
        disabled=_new_logo_multi is None,
        key="btn_multi_replace",
    ):
        _old_refs  = [_mi["bytes"] for _mi in _marked_imgs]
        _new_bytes = _new_logo_multi.read()
        with st.spinner("Replacing logos\u2026"):
            _out_bytes, _count = replace_multiple_logos_any(
                fbytes_sf, name_sf, _old_refs, _new_bytes, _thr_multi
            )
        _tr_count_sf = 0
        if ext_sf in (".pptx", ".docx", ".doc") and _tr_old_sf.strip():
            with st.spinner("Replacing text\u2026"):
                _out_bytes, _tr_count_sf = replace_text_any(
                    _out_bytes, name_sf, _tr_old_sf.strip(), _tr_new_sf
                )
        if _count == 0 and _tr_count_sf == 0:
            st.warning("No replacements made. Try raising the sensitivity slider or checking the text.")
        else:
            _parts = []
            if _count:       _parts.append(f"**{_count}** logo instance(s)")
            if _tr_count_sf: _parts.append(f"**{_tr_count_sf}** text occurrence(s)")
            st.success(f"Replaced {' and '.join(_parts)}.")
        _download_doc("\u2b07\ufe0f Download updated document", _out_bytes, stem_sf, ext_sf)

# ── Step 4: Inspect / region-detect for a single selected image ─────────────

if st.session_state.selected is not None:
    sel  = st.session_state.images[st.session_state.selected]
    name, fbytes = st.session_state.doc_cache
    ext  = os.path.splitext(name.lower())[1]
    stem = os.path.splitext(name)[0]

    st.divider()
    st.subheader("Step 4 — Inspect / region-detect inside one image")

    tab_normal, tab_drill = st.tabs([
        "🔄 Replace this whole image",
        "🔬 Find & replace logo **inside** this image (OCR / template match)",
    ])

    # ── Tab A: Normal whole-image replacement ────────────────────────────────
    with tab_normal:
        c1, c2, c3 = st.columns([1, 1, 1])
        with c1:
            st.markdown("**Selected image**")
            st.image(sel["bytes"], use_container_width=True)
            st.caption("Found in: " + ", ".join(sel["locations"]))
            st.download_button("⬇️ Save this image", to_png(sel["bytes"]),
                               file_name="extracted_image.png", mime="image/png")
        with c2:
            st.markdown("**New logo to replace with**")
            new_logo_whole = st.file_uploader(
                "New logo", type=["png","jpg","jpeg","bmp","webp"],
                label_visibility="collapsed", key="whole_new",
            )
            if new_logo_whole:
                st.image(new_logo_whole, use_container_width=True)
        with c3:
            st.markdown("**Settings**")
            thr_whole = st.slider("Match sensitivity", 0, 20, 5, key="whole_thresh")
            st.caption("Lower = exact match only.")

        if st.button("🔄 Replace whole image", type="primary",
                     disabled=new_logo_whole is None, key="btn_whole"):
            old_bytes = to_png(sel["bytes"])
            new_bytes = new_logo_whole.read()
            with st.spinner("Replacing…"):
                out_bytes, count = replace_logo_any(fbytes, name, old_bytes, new_bytes, thr_whole)
            if count == 0:
                st.warning("No replacements made. Try raising match sensitivity.")
            else:
                st.success(f"Replaced {count} instance(s).")
            _download_doc("⬇️ Download updated document", out_bytes, stem, ext)

    # ── Tab B: Intra-image logo detection + region replace ───────────────────
    with tab_drill:
        st.markdown(
            "Use this when the logo is **baked into a larger background image** "
            "and can't be selected separately in the document editor."
        )
        st.image(sel["bytes"], caption="Selected image — logo will be hunted inside this", use_container_width=True)
        st.divider()

        detect_method = st.radio(
            "Detection method",
            ["OCR — find by text  (recommended for LTIMindtree)", "Template match — find by visual crop"],
            horizontal=True,
            key="detect_method",
        )

        # --- OCR branch ---
        if detect_method.startswith("OCR"):
            st.markdown("**OCR settings**")
            search_text     = st.text_input("Text to search for", value="LTIMindtree", key="ocr_text")
            expand_left     = st.slider(
                "Expand left (to include icon beside text, as % of text width)",
                0, 200, 90, key="ocr_expand",
            )
            v_pad           = st.slider("Vertical padding (%)", 0, 50, 15, key="ocr_vpad")

            if st.button("🔍 Detect logo via OCR", key="btn_ocr"):
                with st.spinner("Running OCR… (first run downloads language model ~50 MB)"):
                    result = detect_by_ocr(
                        sel["bytes"], search_text,
                        expand_left_ratio=expand_left / 100,
                        v_padding_ratio=v_pad / 100,
                    )
                if result is None:
                    st.warning(f'Could not find "{search_text}" in the image. '
                               "Try a partial string (e.g. 'LTI' or 'Mindtree') or switch to Template match.")
                else:
                    st.session_state.bbox = result
                    st.success(f'Found "{result["text"]}" — confidence {result["confidence"]:.0%}')

        # --- Template match branch ---
        else:
            st.markdown(
                "📸 **A screenshot is fine** — you do not need the original logo file. "
                "The tool automatically strips surrounding padding and normalises contrast "
                "so screenshots match reliably. Try to capture just the logo area."
            )
            tmpl_file = st.file_uploader(
                "Screenshot or crop of the old logo",
                type=["png","jpg","jpeg","bmp","webp"], key="tmpl_crop",
            )
            if tmpl_file:
                raw_bytes  = tmpl_file.read()
                crop_bytes = autocrop_screenshot(raw_bytes)
                tc1, tc2 = st.columns(2)
                with tc1:
                    st.caption("Your upload")
                    st.image(raw_bytes, use_container_width=True)
                with tc2:
                    st.caption("After auto-crop (what will be matched)")
                    st.image(crop_bytes, use_container_width=True)
            else:
                raw_bytes = None

            min_score = st.slider("Minimum match score (0 = very loose, 1 = exact)", 0.0, 1.0, 0.40,
                                  step=0.05, key="tmpl_score")

            if raw_bytes and st.button("🔍 Detect logo via template matching", key="btn_tmpl"):
                with st.spinner("Running multi-scale template matching…"):
                    result = detect_by_template(sel["bytes"], raw_bytes, min_score=min_score)
                if result is None:
                    st.warning(
                        f"No match found above score {min_score:.2f}. "
                        "Try lowering the minimum score, or switch to OCR mode."
                    )
                else:
                    st.session_state.bbox = result
                    st.success(f'Match found — score {result["score"]:.0%}')

        # --- Show detected region & replace ---
        if st.session_state.bbox:
            bbox = st.session_state.bbox
            x1, y1, x2, y2 = bbox["x1"], bbox["y1"], bbox["x2"], bbox["y2"]

            st.divider()
            st.markdown("**Detected region** (red box)")
            preview = annotate_bbox(sel["bytes"], x1, y1, x2, y2)
            st.image(preview, use_container_width=True)
            st.caption(f"Region: ({x1}, {y1}) → ({x2}, {y2})  |  size: {x2-x1} × {y2-y1} px")

            # Manual fine-tune
            with st.expander("Fine-tune bounding box"):
                fc1, fc2 = st.columns(2)
                with fc1:
                    x1 = st.number_input("x1 (left)",  value=x1, step=1, key="bb_x1")
                    y1 = st.number_input("y1 (top)",   value=y1, step=1, key="bb_y1")
                with fc2:
                    x2 = st.number_input("x2 (right)", value=x2, step=1, key="bb_x2")
                    y2 = st.number_input("y2 (bottom)", value=y2, step=1, key="bb_y2")
                if st.button("↺ Update preview", key="btn_preview"):
                    preview2 = annotate_bbox(sel["bytes"], x1, y1, x2, y2)
                    st.image(preview2, use_container_width=True)

            st.divider()
            new_logo_drill = st.file_uploader(
                "Upload new logo to place in this region",
                type=["png","jpg","jpeg","bmp","webp"], key="drill_new",
            )
            if new_logo_drill:
                st.image(new_logo_drill, caption="New logo (will be scaled to fit red box)", width=200)

            if st.button("✂️ Replace region & rebuild document", type="primary",
                         disabled=new_logo_drill is None, key="btn_drill_replace"):
                with st.spinner("Replacing region…"):
                    # 1. paint new logo into the selected image
                    modified_img = replace_region_in_image(
                        sel["bytes"], x1, y1, x2, y2, new_logo_drill.read()
                    )
                    # 2. put modified image back into the document
                    out_bytes = embed_image_in_document(fbytes, name, sel, modified_img)

                st.success("Done! The logo region has been replaced in the document.")
                _download_doc("⬇️ Download updated document", out_bytes, stem, ext)

# ── Text Replacement ─────────────────────────────────────────────────────────

if doc_file and not is_zip:
    _name_tr, _fbytes_tr = st.session_state.doc_cache
    _ext_tr = os.path.splitext(_name_tr.lower())[1]
    if _ext_tr in (".pptx", ".docx", ".doc"):
        st.divider()
        st.subheader("🔤 Text Replacement")
        st.caption(
            "Finds and replaces text anywhere in the document — slides, masters, layouts, notes pages. "
            "Works even when the text is split across multiple formatting runs in the XML."
        )
        _tr_c1, _tr_c2 = st.columns(2)
        with _tr_c1:
            _tr_old = st.text_input(
                "Text to find",
                value="LTIMindtree | Privileged and Confidential",
                key="tr_old",
            )
        with _tr_c2:
            _tr_new = st.text_input(
                "Replace with",
                value="LTM | Privileged and Confidential",
                key="tr_new",
            )

        # Offer to also apply logo replacement if logos are already marked + new logo uploaded
        _tr_logo_file = st.session_state.get("multi_new_logo")
        _tr_can_combine = bool(
            st.session_state.get("selected_set")
            and st.session_state.images
            and _tr_logo_file is not None
        )
        _tr_also_logos = False
        if _tr_can_combine:
            _tr_also_logos = st.checkbox(
                "Also apply the logo replacement marked in Step 3 (produces a single combined file)",
                value=True, key="tr_also_logos",
            )

        if st.button("\U0001f524 Replace Text", type="primary", disabled=not _tr_old, key="btn_tr"):
            _working = _fbytes_tr
            _logo_n  = 0
            if _tr_also_logos and _tr_logo_file is not None:
                _tr_logo_file.seek(0)
                _tr_marked = [
                    st.session_state.images[i]
                    for i in sorted(st.session_state.selected_set)
                ]
                _old_refs_tr = [_m["bytes"] for _m in _tr_marked]
                _thr_tr = st.session_state.get("multi_thresh", 8)
                with st.spinner("Replacing logos\u2026"):
                    _working, _logo_n = replace_multiple_logos_any(
                        _working, _name_tr, _old_refs_tr, _tr_logo_file.read(), _thr_tr
                    )
            with st.spinner("Replacing text\u2026"):
                _tr_out, _tr_count = replace_text_any(_working, _name_tr, _tr_old, _tr_new)
            if _tr_count == 0 and _logo_n == 0:
                st.warning(
                    "No replacements made. "
                    "Check the exact spelling and capitalisation, then try again."
                )
            else:
                _parts = []
                if _logo_n:   _parts.append(f"**{_logo_n}** logo instance(s)")
                if _tr_count: _parts.append(f"**{_tr_count}** text occurrence(s)")
                st.success(f"Replaced {' and '.join(_parts)}.")
            _stem_tr = os.path.splitext(_name_tr)[0]
            _download_doc("\u2b07\ufe0f Download updated document", _tr_out, _stem_tr, _ext_tr)

elif doc_file and not is_zip and not st.session_state.images:
    st.info("Click **Scan document for all embedded images** above to get started.")